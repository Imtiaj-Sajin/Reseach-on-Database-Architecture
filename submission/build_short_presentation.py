"""
Build the 12-minute cut of the paper presentation.

Takes the existing DTMS_Paper_Presentation.pptx as-is, keeps slides 1-6
exactly as they stand (including hand edits), drops everything after, and
appends six condensed slides ending on a do / do-not recommendation page.

Every number is measured. Sources: research/results, compare_dbms.py,
paper/paper.tex.
"""
import copy
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = r"g:\codes\Ass\DTMS"
FIG = os.path.join(ROOT, "research", "results", "figures")
SRC = os.path.join(ROOT, "submission", "DTMS_Paper_Presentation.pptx")
OUT = os.environ.get(
    "DECK_OUT",
    os.path.join(ROOT, "submission", "DTMS_Paper_Presentation_Short.pptx"))
KEEP = 6                      # slides 1..6 survive untouched

# ---------------------------------------------------------------- palette
INK    = RGBColor(0x1A, 0x1D, 0x21)
MUTED  = RGBColor(0x70, 0x75, 0x7A)
FAINT  = RGBColor(0xA8, 0xAC, 0xB0)
RULE   = RGBColor(0xD9, 0xD5, 0xCF)
RED    = RGBColor(0xA8, 0x32, 0x2A)
RED_D  = RGBColor(0xE2, 0x76, 0x60)
BLUE   = RGBColor(0x1F, 0x4E, 0x5F)
GREEN  = RGBColor(0x1E, 0x5B, 0x3A)
DARK   = RGBColor(0x14, 0x17, 0x1A)
PAPER  = RGBColor(0xFF, 0xFF, 0xFF)
ONDARK = RGBColor(0xF2, 0xEF, 0xEA)
DIMDK  = RGBColor(0x92, 0x98, 0x9E)

SERIF, SANS, MONO = "Georgia", "Segoe UI", "Consolas"

# ---------------------------------------------------------------- open
prs = Presentation(SRC)
W, H = prs.slide_width, prs.slide_height
L = Inches(0.9)
R = W - Inches(0.9)
CW = R - L

sldIdLst = prs.slides._sldIdLst
for sld in list(sldIdLst)[KEEP:]:
    prs.part.drop_rel(sld.rId)
    sldIdLst.remove(sld)
print("kept %d slides from the existing deck" % len(sldIdLst))

BLANK = prs.slide_layouts[6]

# PowerPoint rewrote the notes master without placeholders when the deck was
# last saved, so a newly added slide gets an empty notes slide and
# notes_text_frame comes back None. Borrow a working body placeholder from a
# slide that already has one.
_NOTES_PH = None
for _sld in prs.slides:
    if _sld.has_notes_slide and _sld.notes_slide.notes_text_frame is not None:
        for _ph in _sld.notes_slide.placeholders:
            if _ph.placeholder_format.type == 2:      # BODY
                _NOTES_PH = copy.deepcopy(_ph._element)
                break
    if _NOTES_PH is not None:
        break
print("notes placeholder template:", "found" if _NOTES_PH is not None else "MISSING")


# ---------------------------------------------------------------- helpers
def slide(dark=False, notes=""):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK if dark else PAPER
    if notes:
        ns = s.notes_slide
        if ns.notes_text_frame is None and _NOTES_PH is not None:
            ns.shapes._spTree.append(copy.deepcopy(_NOTES_PH))
        if ns.notes_text_frame is not None:
            ns.notes_text_frame.text = notes
    s._dark = dark
    return s


def text(s, x, y, w, h, runs, size=16, font=SANS, color=None, bold=False,
         align=PP_ALIGN.LEFT, spacing=1.25, anchor=MSO_ANCHOR.TOP):
    if color is None:
        color = ONDARK if getattr(s, "_dark", False) else INK
    tf = s.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        paras = [[(runs, {})]]
    elif runs and isinstance(runs[0], tuple):
        paras = [runs]
    else:
        paras = runs
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        for t, ov in para:
            r = p.add_run()
            r.text = t
            f = r.font
            f.name = ov.get("font", font)
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.color.rgb = ov.get("color", color)


def line(s, x, y, w, color=None, weight=0.75):
    if color is None:
        color = RULE if not getattr(s, "_dark", False) else RGBColor(0x33, 0x38, 0x3D)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(weight))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False


def eyebrow(s, label, color=None):
    if color is None:
        color = RED if not getattr(s, "_dark", False) else DIMDK
    text(s, L, Inches(0.62), CW, Inches(0.3), label.upper(), size=10.5,
         bold=True, color=color, spacing=1.0)


def heading(s, title, size=30, color=None):
    text(s, L, Inches(1.0), CW, Inches(1.1), title, size=size, font=SERIF,
         color=color, spacing=1.12)


def caption(s, txt):
    text(s, L, H - Inches(0.82), CW, Inches(0.4), txt, size=10.5,
         color=MUTED if not getattr(s, "_dark", False) else DIMDK, spacing=1.2)


def grid(s, x, y, cols, rows, widths, aligns=None, header_size=11,
         cell_size=14, row_h=Inches(0.44), hilite=None, hi_color=None,
         header_h=Inches(0.33)):
    if aligns is None:
        aligns = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * (len(cols) - 1)
    if hi_color is None:
        hi_color = RED
    dark = getattr(s, "_dark", False)
    base = ONDARK if dark else INK
    lab = DIMDK if dark else MUTED
    PAD = Inches(0.10)
    xs, acc = [], x
    for wd in widths:
        xs.append(acc)
        acc += wd
    ws = [wd - 2 * PAD for wd in widths]
    xs = [cx + PAD for cx in xs]
    for cx, cwid, ctxt, al in zip(xs, ws, cols, aligns):
        text(s, cx, y, cwid, header_h, ctxt, size=header_size, bold=True,
             color=lab, align=al, spacing=1.1, anchor=MSO_ANCHOR.BOTTOM)
    ry = y + header_h
    line(s, x, ry, sum(widths),
         color=(base if not dark else RGBColor(0x4A, 0x50, 0x56)), weight=1.0)
    ry += Inches(0.1)
    for i, row in enumerate(rows):
        strong = hilite is not None and i in hilite
        for cx, cwid, ctxt, al in zip(xs, ws, row, aligns):
            text(s, cx, ry + Inches(0.06), cwid, Inches(0.34), str(ctxt),
                 size=cell_size, bold=strong,
                 color=(hi_color if strong else base), align=al, spacing=1.0)
        ry += row_h
        if i < len(rows) - 1:
            line(s, x, ry - Inches(0.04), sum(widths),
                 color=(RULE if not dark else RGBColor(0x2B, 0x30, 0x35)),
                 weight=0.6)
    line(s, x, ry - Inches(0.02), sum(widths),
         color=(base if not dark else RGBColor(0x4A, 0x50, 0x56)), weight=1.0)


def figure(s, name, x, y, max_w, max_h):
    from PIL import Image
    path = os.path.join(FIG, name)
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = max_w, int(max_w / ar)
    if h > max_h:
        h, w = max_h, int(max_h * ar)
    s.shapes.add_picture(path, x + int((max_w - w) / 2), y, w, h)


def bignum(s, value, label, x, y, w, size=76, color=None, lab_color=None):
    dark = getattr(s, "_dark", False)
    if color is None:
        color = ONDARK if dark else INK
    if lab_color is None:
        lab_color = DIMDK if dark else MUTED
    text(s, x, y, w, Inches(1.4), value, size=size, font=SERIF, color=color,
         spacing=0.95)
    text(s, x, y + Inches(size / 72.0 * 1.02) + Inches(0.13), w, Inches(0.9),
         label, size=13, color=lab_color, spacing=1.25)


def pagenum(s, n):
    text(s, R - Inches(0.7), H - Inches(0.58), Inches(0.7), Inches(0.3), str(n),
         size=10, align=PP_ALIGN.RIGHT, spacing=1.0,
         color=FAINT if not getattr(s, "_dark", False)
         else RGBColor(0x4A, 0x50, 0x56))


# ================================================================ 7 REPERTOIRE
s = slide(notes=(
    "So why is MariaDB wrong so much more often? It is not that it costs "
    "plans less carefully. It is that it has fewer plans to choose from. "
    "This counts what each engine actually picked when handed every index. "
    "PostgreSQL answers almost everything with a bitmap scan, which sorts "
    "the matching row locations into disk order before fetching them. "
    "MariaDB never forms one. It either does a direct index lookup or gives "
    "up and scans the whole table, sixty six times, where PostgreSQL did it "
    "nine times. That single difference explains the gap."))
eyebrow(s, "Why the two engines differ")
heading(s, "The gap is the plans available, not the care taken")
grid(s, L, Inches(2.25),
     ["What the engine chose", "PostgreSQL", "MariaDB"],
     [["Bitmap scan  (sorts rows into disk order)", "346", "7"],
      ["Direct index lookup", "19", "301"],
      ["Gave up, scanned the whole table", "9", "66"]],
     [Inches(5.6), Inches(2.0), Inches(1.8)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
     cell_size=15, row_h=Inches(0.55))
line(s, L, Inches(4.35), Inches(9.4))
text(s, L, Inches(4.75), Inches(6.2), Inches(2.0), [
    [("A bitmap scan puts a ceiling on how bad a mistake can be, because "
      "even a wrong choice still reads the disk in order.",
      {"size": 16.5})],
    [("", {"size": 8})],
    [("MariaDB has no ceiling. A mistaken index scan becomes random "
      "single-row fetches.", {"size": 16.5, "bold": True, "color": RED})],
], spacing=1.35)
text(s, Inches(7.6), Inches(4.75), Inches(4.8), Inches(2.0), [
    [("This is one design decision,", {"size": 15, "color": MUTED})],
    [("not a tuning gap. It is why", {"size": 15, "color": MUTED})],
    [("PostgreSQL errs rarely and", {"size": 15, "color": MUTED})],
    [("mildly, and MariaDB errs", {"size": 15, "color": MUTED})],
    [("often and sometimes badly.", {"size": 15, "color": MUTED})],
], spacing=1.3)
caption(s, "Counted over the free-choice configuration, all four query families, 1,000,000 rows.")
pagenum(s, 7)

# ================================================================ 8 PHENOMENON
s = slide(notes=(
    "Everything from here is about PostgreSQL alone, and I want to be clear "
    "about that because MariaDB does not have a BRIN index at all, so it "
    "cannot make this particular mistake. "
    "On slide six I said PostgreSQL chooses well, 2.3 percent bad. That "
    "number was measured without a BRIN index. Now watch what happens when "
    "I add one. Same engine, same data, same queries. The only difference "
    "between these two rows is whether a BRIN index exists on the column "
    "beside the B-tree. Bad choices go from 2.3 percent to 73.3 percent, "
    "and the worst case goes from 1.4 times to nearly 19. So this is not a "
    "contradiction of slide six. It is what adding one supposedly free "
    "index does to it."))
eyebrow(s, "PostgreSQL only, from here to slide 10")
heading(s, "Adding one more index broke it")
text(s, L, Inches(2.05), Inches(11.4), Inches(0.5),
     "Same engine, same data, same queries. The only difference is whether "
     "a BRIN index exists.", size=17, color=INK, spacing=1.25)
grid(s, L, Inches(2.85),
     ["Indexes available to the planner", "Chose badly", "Worst case"],
     [["B-tree only  (this is the 2.3% from slide 6)", "2.3%", "1.4x"],
      ["B-tree + a BRIN index on the same column", "73.3%", "18.8x"]],
     [Inches(6.2), Inches(1.9), Inches(1.7)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
     cell_size=15.5, row_h=Inches(0.62), hilite=[1])
line(s, L, Inches(4.85), Inches(9.8))
text(s, L, Inches(5.2), Inches(6.0), Inches(1.5), [
    [("People add BRIN indexes because they are tiny and look free.",
      {"size": 16})],
    [("", {"size": 8})],
    [("It is free in disk space. It is not free in plan quality.",
      {"size": 16, "bold": True, "color": RED})],
], spacing=1.32)
text(s, Inches(7.4), Inches(5.2), Inches(5.0), Inches(1.5), [
    [("MariaDB has no BRIN index,", {"size": 15, "color": MUTED})],
    [("so it cannot make this mistake", {"size": 15, "color": MUTED})],
    [("at all. That is how we know it", {"size": 15, "color": MUTED})],
    [("is a PostgreSQL costing defect,", {"size": 15, "color": MUTED})],
    [("not something unavoidable.", {"size": 15, "color": MUTED})],
], spacing=1.28)
caption(s, "1,000,000 rows. The effect appears only while the table still fits in the 128 MB buffer pool.")
pagenum(s, 8)

# ================================================================ 9 NOT ESTIMATES
s = slide(notes=(
    "So why did it choose badly? The textbook answer is that the database "
    "misjudged how many rows would come back. I checked. Of those 113 bad "
    "choices, 98 percent had the row count essentially right. The median "
    "estimation error was 1.017, where a perfect estimate is 1.000. "
    "So it knew how many rows were coming. It knew both indexes existed. "
    "And it still picked the slower one. That rules out the usual "
    "explanation, and it rules out the usual fix, because running ANALYZE "
    "or adding statistics cannot improve a number that was already correct. "
    "Which leaves one question: if the estimates were right, what did it "
    "get wrong?"))
eyebrow(s, "Ruling out the obvious cause")
heading(s, "It was not a case of bad estimates")
bignum(s, "98.2%",
       "of those 113 bad choices had the row count\n"
       "essentially correct  (median error 1.017)",
       L, Inches(2.3), Inches(6.3), size=84, color=RED)
line(s, L, Inches(5.0), Inches(6.3))
text(s, L, Inches(5.35), Inches(6.3), Inches(1.3), [
    [("The 113 bad choices from the previous slide, PostgreSQL at 1,000,000 "
      "rows. A perfect estimate is 1.000.", {"size": 14.5, "color": MUTED})],
], spacing=1.3)
text(s, Inches(7.7), Inches(2.4), Inches(4.7), Inches(3.9), [
    [("It knew how many rows were coming back.", {"size": 17, "bold": True})],
    [("It knew both indexes existed.", {"size": 17, "bold": True})],
    [("It still picked the slower one.", {"size": 17, "bold": True, "color": RED})],
    [("", {"size": 14})],
    [("So running ANALYZE or adding statistics cannot fix this. You cannot "
      "improve a number that was already right.", {"size": 15, "color": MUTED})],
    [("", {"size": 10})],
    [("Which leaves the question the next slide answers: if the estimates "
      "were right, what did it get wrong?", {"size": 15.5, "bold": True})],
], spacing=1.3)
pagenum(s, 9)

# ================================================================ 10 CAUSE
s = slide(dark=True, notes=(
    "Here is the answer, and I found it by reading the PostgreSQL source "
    "code. There is a function called choose_bitmap_and. When two indexes "
    "can both answer the same query, it keeps one and throws the other "
    "away, and it decides by asking which index is cheaper to read. That "
    "sounds reasonable. The problem is that reading the index is only half "
    "the job. After reading it, the query still has to go and fetch the "
    "actual rows from the table, and that cost is not counted in the "
    "comparison. A BRIN index is tiny, so it is always cheaper to read, so "
    "it always wins. But BRIN only narrows the search to a block range, so "
    "the query then has to scan far more rows. Look at the numbers: the "
    "index that costs 17 times less to read is the one that runs 194 times "
    "slower. It compares the wrong quantity."))
eyebrow(s, "The actual cause")
heading(s, "It compares the wrong thing", color=ONDARK)
text(s, L, Inches(2.0), Inches(11.4), Inches(0.9), [
    [("When two indexes can answer the same query, PostgreSQL keeps "
      "whichever is ", {"size": 17, "color": ONDARK}),
     ("cheaper to read", {"size": 17, "bold": True, "color": RED_D}),
     (", and ignores the cost of fetching the rows afterwards.",
      {"size": 17, "color": ONDARK})],
], spacing=1.3)
grid(s, L, Inches(3.2),
     ["Index", "Cost it used to decide", "How long the query actually took"],
     [["BRIN  (kept)", "12.13", "65.95 ms"],
      ["B-tree  (discarded)", "213.35", "0.34 ms"]],
     [Inches(2.9), Inches(3.0), Inches(4.0)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
     cell_size=16, row_h=Inches(0.62))
line(s, L, Inches(5.25), Inches(9.9))
text(s, L, Inches(5.6), Inches(6.4), Inches(1.2), [
    [("The index that costs 17x less to read is the one that runs "
      "194x slower.", {"size": 18, "font": SERIF, "color": RED_D})],
], spacing=1.25)
text(s, Inches(7.9), Inches(5.55), Inches(4.5), Inches(1.3), [
    [("A BRIN index is tiny, so it always wins that comparison. But it only "
      "narrows the search to a range of blocks, so the query then has far "
      "more rows to check.", {"size": 14, "color": DIMDK})],
], spacing=1.28)
caption(s, "choose_bitmap_and() in optimizer/path/indxpath.c. A design consequence, which is why no setting turns it off.")
pagenum(s, 10)

# ================================================================ 11 OPPOSITE
s = slide(notes=(
    "Measuring both engines across all seven table sizes shows they do not "
    "just differ by degree, they fail in opposite ways. As tables grow, "
    "PostgreSQL gets wrong more often but its worst case stays small, never "
    "past 3.7 times. MariaDB gets wrong less often but its worst case "
    "climbs to 31 times. The practical point is the one on the right: a "
    "database that is wrong often but never badly is easier to live with "
    "than one wrong rarely but catastrophically. If you summarise either "
    "engine with a single number you lose exactly this."))
eyebrow(s, "How they fail as data grows")
heading(s, "The two engines fail in opposite directions")
figure(s, "fig8_degradation_shape.png", L, Inches(2.0), Inches(6.6), Inches(3.3))
text(s, L, Inches(5.5), Inches(6.6), Inches(1.2),
     "Left: how often it errs.   Right: how badly, at worst.",
     size=13, color=MUTED, spacing=1.25)
text(s, Inches(8.0), Inches(2.1), Inches(4.4), Inches(4.4), [
    [("PostgreSQL", {"size": 16, "bold": True, "color": BLUE})],
    [("gets wrong more often as tables grow, 0% to 36%, but never by more "
      "than 3.7x.", {"size": 15})],
    [("", {"size": 10})],
    [("MariaDB", {"size": 16, "bold": True, "color": RED})],
    [("gets wrong less often, 53% down to 12%, but its worst case climbs to "
      "30.9x.", {"size": 15})],
    [("", {"size": 12})],
    [("Wrong often but never badly is easier to run than wrong rarely but "
      "catastrophically.", {"size": 15.5, "bold": True})],
], spacing=1.32)
caption(s, "Seven table scales, 1 to 10 million rows. Each engine judged only against plans the other could also have formed.")
pagenum(s, 11)

# ================================================================ 12 PRACTICES
s = slide(notes=(
    "Pulling it together. We tested four pieces of advice that appear in "
    "official documentation and well-regarded blogs. Every one of them is "
    "sensible. Every one of them, under conditions we can now name, makes "
    "things worse. And look at the middle column: in three of the four the "
    "advice does exactly what it promises. Extended statistics really do "
    "fix the estimate, by a factor of 108. The query still gets slower. "
    "That is the pattern of the whole study: the advice improves a proxy, "
    "and the proxy is not what determines the outcome."))
eyebrow(s, "The pattern")
heading(s, "Four recommended practices, all four backfire")
line(s, L, Inches(2.05), CW)
text(s, Inches(4.7), Inches(2.15), Inches(3.6), Inches(0.3),
     "WHAT IT PROMISES", size=10, bold=True, color=MUTED, spacing=1.0)
text(s, Inches(8.6), Inches(2.15), Inches(3.8), Inches(0.3),
     "WHAT WE MEASURED", size=10, bold=True, color=RED, spacing=1.0)
rows = [
    ("Extended statistics", "for correlated columns",
     "Fixes the estimate 108x", "Query up to 2.7x slower"),
    ("Lower random_page_cost", "because the disk is an SSD",
     "Default is already near optimal", "Lowering it 2.3x worse"),
    ("Add a BRIN index", "it is small and cheap",
     "Costs almost no space", "Bad choices 2.3% to 73.3%"),
    ("Enable Multi-Range Read", "MariaDB, ships switched off",
     "The mechanism does engage", "Workload 7% slower"),
]
y = Inches(2.6)
for name, why, does, costs in rows:
    text(s, L, y, Inches(3.6), Inches(0.8), [
        [(name, {"size": 15.5, "bold": True})],
        [(why, {"size": 12.5, "color": MUTED})],
    ], spacing=1.3)
    text(s, Inches(4.7), y + Inches(0.06), Inches(3.6), Inches(0.6), does,
         size=14, color=BLUE, spacing=1.25)
    text(s, Inches(8.6), y + Inches(0.06), Inches(3.8), Inches(0.6), costs,
         size=14, color=RED, bold=True, spacing=1.25)
    y += Inches(1.0)
    if y < Inches(6.2):
        line(s, L, y - Inches(0.15), CW)
text(s, L, Inches(6.5), CW, Inches(0.5),
     "In three of four the advice delivers exactly what it promises, and the "
     "plan gets worse anyway.", size=15, bold=True, spacing=1.25)
pagenum(s, 12)

# ================================================================ 13 RECOMMEND
s = slide(notes=(
    "So what do we actually recommend. On the left, four things to do. "
    "Measure the whole workload, not one query, because we had a 24 times "
    "effect on one query that vanished across the workload. Test more than "
    "two table sizes, because the problem moved in two directions at once "
    "and two points would have given us the opposite answer. Report how "
    "often and how badly separately, because they move in opposite "
    "directions. And keep the shipped defaults unless your own measurement "
    "says otherwise. On the right, four things not to do, each with the "
    "number that justifies it. And the line at the bottom is the one "
    "sentence I would want people to remember: every one of these was good "
    "advice, it just never came with the conditions attached."))
eyebrow(s, "What we recommend")
heading(s, "What to do, and what not to do")
line(s, L, Inches(2.0), CW)
text(s, L, Inches(2.35), Inches(5.8), Inches(0.35),
     "DO", size=13, bold=True, color=GREEN, spacing=1.0)
do = [
    ("Measure the whole workload, not one query",
     "A 24x tuning effect vanished across the workload."),
    ("Test more than two table sizes",
     "Frequency and severity moved in opposite directions."),
    ("Report how often and how badly, separately",
     "One number hides which engine is safer to run."),
    ("Keep shipped defaults until measured otherwise",
     "The default random_page_cost was within 6% of optimal."),
]
y = Inches(2.8)
for h, b in do:
    text(s, L, y, Inches(0.3), Inches(0.4), "+", size=15, bold=True,
         color=GREEN, spacing=1.0)
    text(s, L + Inches(0.38), y - Inches(0.02), Inches(5.4), Inches(0.9), [
        [(h, {"size": 14.5, "bold": True})],
        [(b, {"size": 12.5, "color": MUTED})],
    ], spacing=1.28)
    y += Inches(0.92)

text(s, Inches(7.2), Inches(2.35), Inches(5.2), Inches(0.35),
     "DO NOT", size=13, bold=True, color=RED, spacing=1.0)
dont = [
    ("Add a BRIN index next to a B-tree",
     "Same column, table in memory: bad choices 2.3% to 73.3%, worst 194x."),
    ("Add extended statistics for a faster plan",
     "Estimate improved 108x, query 1.3x to 2.7x slower."),
    ("Lower random_page_cost for an SSD",
     "At 1.0 the workload was 2.3x worse than the default."),
    ("Enable MariaDB MRR at its default buffer",
     "7% slower; helps only at 8 MB and only on narrow queries."),
]
y = Inches(2.8)
for h, b in dont:
    text(s, Inches(7.2), y, Inches(0.3), Inches(0.4), "\u2013", size=15,
         bold=True, color=RED, spacing=1.0)
    text(s, Inches(7.58), y - Inches(0.02), Inches(4.85), Inches(0.9), [
        [(h, {"size": 14.5, "bold": True})],
        [(b, {"size": 12.5, "color": MUTED})],
    ], spacing=1.28)
    y += Inches(0.92)

line(s, L, Inches(6.5), CW)
text(s, L, Inches(6.75), Inches(10.9), Inches(0.6),
     "All four were reasonable advice. None of it shipped with the "
     "conditions under which it stops being true.",
     size=15, font=SERIF, bold=True, spacing=1.2)
pagenum(s, 13)

prs.save(OUT)
print("wrote", OUT)
print("total slides:", len(prs.slides._sldIdLst))
