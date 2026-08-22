"""
Build the viva presentation deck.

Every number in this deck is copied from research/results and STATUS.md.
Nothing here is invented. Figures are the same PNGs used in the report.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = r"g:\codes\Ass\DTMS"
FIG = os.path.join(ROOT, "research", "results", "figures")
OUT = os.path.join(ROOT, "submission", "DTMS_Presentation.pptx")

# ---------------------------------------------------------------- palette
INK    = RGBColor(0x1A, 0x1D, 0x21)
PAPER  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x70, 0x75, 0x7A)
FAINT  = RGBColor(0xA8, 0xAC, 0xB0)
RULE   = RGBColor(0xD9, 0xD5, 0xCF)
RED    = RGBColor(0xA8, 0x32, 0x2A)
RED_D  = RGBColor(0xE2, 0x76, 0x60)     # red legible on the dark slides
BLUE   = RGBColor(0x1F, 0x4E, 0x5F)
DARK   = RGBColor(0x14, 0x17, 0x1A)
ONDARK = RGBColor(0xF2, 0xEF, 0xEA)
DIMDK  = RGBColor(0x92, 0x98, 0x9E)

SERIF = "Georgia"
SANS  = "Segoe UI"
MONO  = "Consolas"

W, H = Inches(13.333), Inches(7.5)
L = Inches(0.9)                 # left margin
R = W - Inches(0.9)             # right edge
CW = R - L                      # content width

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide(dark=False, notes=""):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK if dark else PAPER
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    s._dark = dark
    return s


def text(s, x, y, w, h, runs, size=16, font=SANS, color=None, bold=False,
         align=PP_ALIGN.LEFT, spacing=1.25, anchor=MSO_ANCHOR.TOP, space_after=0):
    """runs: str, or list of (text, {overrides}) tuples, or list of such lists
    (one inner list per paragraph)."""
    if color is None:
        color = ONDARK if getattr(s, "_dark", False) else INK
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    if isinstance(runs, str):
        paras = [[(runs, {})]]
    elif runs and isinstance(runs[0], tuple):
        paras = [runs]
    else:
        paras = runs

    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        for t, ov in para:
            r = p.add_run()
            r.text = t
            f = r.font
            f.name = ov.get("font", font)
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.color.rgb = ov.get("color", color)
    return box


def line(s, x, y, w, color=None, weight=0.75):
    if color is None:
        color = RULE if not getattr(s, "_dark", False) else RGBColor(0x33, 0x38, 0x3D)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(weight))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def eyebrow(s, label, color=None):
    if color is None:
        color = RED if not getattr(s, "_dark", False) else DIMDK
    text(s, L, Inches(0.62), CW, Inches(0.3),
         label.upper(), size=10.5, font=SANS, bold=True, color=color, spacing=1.0)


def heading(s, title, y=Inches(1.0), size=29, color=None, w=None):
    text(s, L, y, w or CW, Inches(1.1), title, size=size, font=SERIF,
         color=color, spacing=1.12)


def caption(s, txt, y=None):
    text(s, L, y or (H - Inches(0.82)), CW, Inches(0.4), txt,
         size=10.5, color=MUTED if not getattr(s, "_dark", False) else DIMDK,
         spacing=1.2)


def grid(s, x, y, cols, rows, widths, aligns=None, header_size=11,
         cell_size=13.5, row_h=Inches(0.42), hilite=None, hi_color=None,
         header_h=Inches(0.33)):
    """Editorial table: no fills, no boxes, hairline rules only."""
    if aligns is None:
        aligns = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * (len(cols) - 1)
    if hi_color is None:
        hi_color = RED
    dark = getattr(s, "_dark", False)
    base = ONDARK if dark else INK
    lab = DIMDK if dark else MUTED

    PAD = Inches(0.10)          # gutter, so adjacent columns never touch
    xs, acc = [], x
    for wd in widths:
        xs.append(acc)
        acc += wd
    ws = [wd - 2 * PAD for wd in widths]
    xs = [cx + PAD for cx in xs]

    # header
    for cx, cwid, ctxt, al in zip(xs, ws, cols, aligns):
        text(s, cx, y, cwid, header_h, ctxt, size=header_size, font=SANS,
             bold=True, color=lab, align=al, spacing=1.1,
             anchor=MSO_ANCHOR.BOTTOM)
    ry = y + header_h
    line(s, x, ry, sum(widths), color=(base if not dark else RGBColor(0x4A, 0x50, 0x56)),
         weight=1.0)
    ry += Inches(0.1)

    for i, row in enumerate(rows):
        strong = hilite is not None and i in hilite
        for cx, cwid, ctxt, al in zip(xs, ws, row, aligns):
            col = hi_color if strong else base
            text(s, cx, ry + Inches(0.06), cwid, Inches(0.34), str(ctxt),
                 size=cell_size, font=SANS, bold=strong, color=col,
                 align=al, spacing=1.0)
        ry += row_h
        if i < len(rows) - 1:
            line(s, x, ry - Inches(0.04), sum(widths),
                 color=(RULE if not dark else RGBColor(0x2B, 0x30, 0x35)), weight=0.6)
    line(s, x, ry - Inches(0.02), sum(widths),
         color=(base if not dark else RGBColor(0x4A, 0x50, 0x56)), weight=1.0)
    return ry


def figure(s, name, x, y, max_w, max_h):
    from PIL import Image
    path = os.path.join(FIG, name)
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = max_w
    h = int(w / ar)
    if h > max_h:
        h = max_h
        w = int(h * ar)
    s.shapes.add_picture(path, x + int((max_w - w) / 2), y, w, h)
    return h


def bignum(s, value, label, x, y, w, size=76, color=None, lab_color=None):
    dark = getattr(s, "_dark", False)
    if color is None:
        color = ONDARK if dark else INK
    if lab_color is None:
        lab_color = DIMDK if dark else MUTED
    text(s, x, y, w, Inches(1.4), value, size=size, font=SERIF, color=color, spacing=0.95)
    text(s, x, y + Inches(size / 72.0 * 1.02) + Inches(0.13), w, Inches(0.8), label,
         size=12.5, font=SANS, color=lab_color, spacing=1.25)


def pagenum(s, n):
    if n == 0:
        return
    text(s, R - Inches(0.7), H - Inches(0.58), Inches(0.7), Inches(0.3), str(n),
         size=10, color=FAINT if not getattr(s, "_dark", False) else RGBColor(0x4A, 0x50, 0x56),
         align=PP_ALIGN.RIGHT, spacing=1.0)


# ================================================================ 1. TITLE
s = slide(dark=True, notes=(
    "Good morning. My project is an experimental study of how PostgreSQL and "
    "MariaDB choose which index to use, and whether the tuning advice people "
    "give for it actually works. Everything I show is measured on my own "
    "machine and the code is public."))
line(s, L, Inches(1.55), Inches(1.4), color=RED, weight=2.5)
text(s, L, Inches(1.85), Inches(10.6), Inches(2.3),
     "Measuring the Impact of Index Types and Tuning Settings on "
     "Query Execution Plans in PostgreSQL and MariaDB",
     size=35, font=SERIF, color=ONDARK, spacing=1.14)
text(s, L, Inches(4.35), Inches(10.6), Inches(0.4),
     "Type 4  |  Experimental and Benchmarking Project  |  Individual Submission",
     size=13, color=DIMDK, spacing=1.2)
line(s, L, Inches(4.95), CW, color=RGBColor(0x33, 0x38, 0x3D))
text(s, L, Inches(5.25), Inches(5.4), Inches(1.3), [
    [("Md. Imtiaj Alam Sajin", {"size": 15, "bold": True, "color": ONDARK})],
    [("26-94090-2   |   26-94090-2@student.aiub.edu", {"size": 12, "color": DIMDK})],
], spacing=1.45)
text(s, Inches(7.2), Inches(5.25), Inches(5.2), Inches(1.3), [
    [("Supervisor", {"size": 11, "color": DIMDK})],
    [("Dr. Ashraf Uddin", {"size": 15, "bold": True, "color": ONDARK})],
], spacing=1.45)
text(s, L, H - Inches(0.85), CW, Inches(0.4),
     "github.com/Imtiaj-Sajin/Reseach-on-Database-Architecture",
     size=11.5, font=MONO, color=DIMDK, spacing=1.2)

# ================================================================ 2. PREMISE
s = slide(notes=(
    "Here is the gap I started from. If you search for how to speed up a slow "
    "query you get thousands of pages of advice. Add this index, change this "
    "setting. It is written by good engineers and it sounds sensible. But "
    "almost none of it has been tested under controlled conditions, so nobody "
    "knows when it works and when it does not. That is what I measured."))
eyebrow(s, "The gap")
heading(s, "Index tuning advice is everywhere.\nAlmost none of it has been measured.", size=32)
line(s, L, Inches(2.85), CW)
text(s, L, Inches(3.2), Inches(6.1), Inches(3.0), [
    [("Every query needs one decision made for it:", {"size": 16})],
    [("scan the whole table, or use an index, and if so, which one.",
      {"size": 16, "bold": True})],
    [("", {"size": 8})],
    [("The same query on the same data can be a hundred times slower "
      "depending on that single choice, and the user never sees it happen.",
      {"size": 16, "color": MUTED})],
], spacing=1.42)
text(s, Inches(7.6), Inches(3.2), Inches(4.8), Inches(3.0), [
    [("So a body of advice grew around it", {"size": 13, "bold": True, "color": RED})],
    [("", {"size": 6})],
    [("Create extended statistics for correlated columns", {"size": 14})],
    [("Lower random_page_cost because you have an SSD", {"size": 14})],
    [("Add a BRIN index, it is tiny and cheap", {"size": 14})],
    [("Enable Multi-Range Read in MariaDB", {"size": 14})],
    [("", {"size": 6})],
    [("All four are reasonable. I tested all four.",
      {"size": 14, "bold": True})],
], spacing=1.5)
pagenum(s, 2)

# ================================================================ 3. QUESTIONS
s = slide(notes=(
    "I turned that into five questions. The first two ask how good the "
    "optimiser actually is and why it fails. The next two test the advice "
    "itself. The last one asks whether what I find is a PostgreSQL quirk or "
    "something deeper, which is why I added a second database."))
eyebrow(s, "Research questions")
heading(s, "Five questions")
line(s, L, Inches(2.05), CW)
qs = [
    ("RQ1", "How often does the planner pick a slower access path than one it already had?"),
    ("RQ2", "Is bad selectivity estimation the cause, as the literature assumes?"),
    ("RQ3", "Do the recommended tuning practices improve the plan, or only the estimate?"),
    ("RQ4", "Where are the boundaries in table size and memory at which behaviour flips?"),
    ("RQ5", "Which findings are general to relational optimisers, and which belong to one product?"),
]
y = Inches(2.4)
for tag, q in qs:
    text(s, L, y, Inches(0.85), Inches(0.4), tag, size=13, font=SANS, bold=True, color=RED, spacing=1.0)
    text(s, L + Inches(0.95), y - Inches(0.03), Inches(10.5), Inches(0.5), q, size=16.5, spacing=1.25)
    y += Inches(0.86)
caption(s, "RQ5 is why the study covers two independently written database engines rather than one.")
pagenum(s, 3)

# ================================================================ 4. METHOD
s = slide(notes=(
    "I did not use a standard benchmark like TPC-H, and that was deliberate. "
    "With a fixed dataset you cannot separate cause from coincidence. I "
    "generate the data myself, so I can move one property at a time and hold "
    "the rest still, and more importantly I know the true answer for every "
    "query. That is what lets me say the database guessed wrong instead of "
    "guessing along with it."))
eyebrow(s, "Method")
heading(s, "Generated data, one property at a time")
text(s, L, Inches(2.0), Inches(5.7), Inches(2.4), [
    [("Why not TPC-H or a real dataset?", {"size": 15, "bold": True})],
    [("", {"size": 6})],
    [("With a fixed dataset you cannot separate cause from coincidence, and "
      "you do not know the correct answer. Generating it means I control one "
      "property at a time and I know the true row count for every query, so I "
      "can say the planner was wrong rather than guess alongside it.",
      {"size": 15, "color": MUTED})],
], spacing=1.4)
grid(s, Inches(7.0), Inches(2.05),
     ["Property varied", "Meaning", "Levels"],
     [["Value skew", "Zipfian, some values dominate", "4"],
      ["Predicate correlation", "Two columns move together", "5"],
      ["Physical clustering", "Rows stored in value order", "4"]],
     [Inches(1.95), Inches(2.55), Inches(0.75)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT],
     cell_size=12.5)
line(s, L, Inches(4.75), CW)
text(s, L, Inches(5.05), CW, Inches(1.4), [
    [("11 dataset configurations", {"size": 15, "bold": True}),
     ("  x  ", {"size": 15, "color": FAINT}),
     ("10 index configurations", {"size": 15, "bold": True}),
     ("  x  ", {"size": 15, "color": FAINT}),
     ("4 query families", {"size": 15, "bold": True}),
     ("  x  ", {"size": 15, "color": FAINT}),
     ("8 selectivity targets", {"size": 15, "bold": True})],
    [("", {"size": 7})],
    [("Every query runs 7 times and the first is discarded. Parallelism and JIT "
      "are switched off, the buffer pool is pinned at 128 MB, and each run is "
      "seeded so it reproduces exactly.", {"size": 14, "color": MUTED})],
], spacing=1.4)
pagenum(s, 4)

# ================================================================ 5. MATRIX
s = slide(notes=(
    "This is the size of the experiment. Seven table sizes on both databases, "
    "from one million rows up to ten million, plus configuration sweeps on "
    "each. Twenty three thousand four hundred and eighty measurements in "
    "total. The runs took several days and survived two power cuts, which is "
    "why every measurement is written to disk as it is taken and the run "
    "resumes where it stopped."))
eyebrow(s, "Scale")
heading(s, "23,480 measurements")
rows = [
    ["1,000,000", "112 MB", "2,398", "1,232"],
    ["1,250,000", "140 MB", "654", "336"],
    ["1,500,000", "168 MB", "654", "336"],
    ["2,000,000", "223 MB", "654", "336"],
    ["3,000,000", "335 MB", "654", "336"],
    ["5,000,000", "558 MB", "654", "336"],
    ["10,000,000", "1,116 MB", "2,398", "1,232"],
    ["Configuration sweeps", "", "6,230", "5,040"],
    ["Total", "", "14,296", "9,184"],
]
grid(s, L, Inches(2.05),
     ["Table size", "Heap", "PostgreSQL 17.1", "MariaDB 10.4.28"],
     rows, [Inches(3.0), Inches(1.5), Inches(2.2), Inches(2.2)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
     row_h=Inches(0.395), cell_size=13, hilite=[8], hi_color=INK)
text(s, Inches(10.0), Inches(2.4), Inches(2.6), Inches(3.4), [
    [("Both engines cover the", {"size": 13.5, "color": MUTED})],
    [("same seven table sizes,", {"size": 13.5, "color": MUTED})],
    [("so the comparison is", {"size": 13.5, "color": MUTED})],
    [("symmetric.", {"size": 13.5, "color": MUTED})],
    [("", {"size": 10})],
    [("Buffer pool held at", {"size": 13.5, "color": MUTED})],
    [("128 MB throughout, so", {"size": 13.5, "color": MUTED})],
    [("the memory boundary", {"size": 13.5, "color": MUTED})],
    [("falls between rows one", {"size": 13.5, "color": MUTED})],
    [("and two.", {"size": 13.5, "color": MUTED})],
], spacing=1.32)
caption(s, "Runs are resumable per measurement. Two power cuts during collection cost no data.")
pagenum(s, 5)

# ================================================================ 6. METRIC
s = slide(notes=(
    "One definition I need before the results. Regret is how many times slower "
    "the plan the database chose was, compared to the best plan that was "
    "available to it at that moment. Not the theoretical best, the best it "
    "could actually have picked. A regret of one means it did the best it "
    "could. This matters because it separates the optimiser being wrong from "
    "the database simply being slow."))
eyebrow(s, "The metric")
heading(s, "Access-path regret")
text(s, L, Inches(2.15), Inches(7.7), Inches(0.7),
     "R  =  time of the plan actually chosen  /  time of the fastest plan available",
     size=16, font=SERIF, color=BLUE, spacing=1.2)
line(s, L, Inches(3.05), Inches(7.7))
text(s, L, Inches(3.35), Inches(7.7), Inches(2.4), [
    [("R = 1.0 means the database made the best choice it could.",
      {"size": 16, "bold": True})],
    [("", {"size": 7})],
    [("The comparison is against plans the engine really had, forced one at a "
      "time by building each index configuration separately and re-running the "
      "same query. So this measures the decision, not the engine's raw speed.",
      {"size": 15, "color": MUTED})],
    [("", {"size": 7})],
    [("Estimation quality is reported separately as q-error, the standard "
      "measure: the factor by which the row estimate misses the truth.",
      {"size": 15, "color": MUTED})],
], spacing=1.4)
bignum(s, "1.0", "best choice available", Inches(9.1), Inches(2.35), Inches(3.3), size=62, color=BLUE)
bignum(s, "194x", "worst case I measured", Inches(9.1), Inches(4.25), Inches(3.3), size=62, color=RED)
pagenum(s, 6)

# ================================================================ 7. F1
s = slide(notes=(
    "First result, and it is the one that frames everything else. On identical "
    "data with identical queries, PostgreSQL picks the best plan available "
    "ninety seven point seven percent of the time. So the failures I show you "
    "later are specific, identifiable problems, not a generally weak "
    "optimiser. MariaDB is wrong far more often, but notice its errors are "
    "milder. The reason is visible in the plans: PostgreSQL has bitmap scans, "
    "which sort matching rows into disk order first. MariaDB has nothing "
    "equivalent, so it either does a direct lookup or gives up and scans the "
    "whole table."))
eyebrow(s, "Finding 1  /  RQ1")
heading(s, "PostgreSQL is good at this. MariaDB much less so.")
grid(s, L, Inches(2.15),
     ["Database", "Queries", "Chose badly", "Median regret", "Worst case"],
     [["PostgreSQL 17.1", "130", "2.3%", "1.01", "1.39"],
      ["MariaDB 10.4.28", "153", "66.0%", "1.34", "7.91"]],
     [Inches(2.9), Inches(1.35), Inches(1.6), Inches(1.7), Inches(1.5)],
     aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 4,
     row_h=Inches(0.48), cell_size=14.5)
line(s, L, Inches(3.85), Inches(9.05))
text(s, L, Inches(4.15), Inches(6.2), Inches(2.4), [
    [("Why the difference is structural, not a tuning gap",
      {"size": 14, "bold": True, "color": RED})],
    [("", {"size": 6})],
    [("PostgreSQL uses bitmap scans, which collect the matching row locations "
      "and sort them into disk order before fetching. MariaDB has no "
      "equivalent enabled, so it either performs a direct index lookup or "
      "falls back to scanning the whole table. It scanned the whole table for "
      "66 queries where PostgreSQL never did once.",
      {"size": 14.5, "color": MUTED})],
], spacing=1.38)
figure(s, "fig2_regret_by_selectivity.png", Inches(7.2), Inches(4.0), Inches(5.2), Inches(2.6))
caption(s, "Left column: measured on 1,000,000 rows, identical data and queries on both engines.")
pagenum(s, 7)

# ================================================================ 8. F2
s = slide(notes=(
    "This is the result I did not expect. The whole query optimisation "
    "literature says bad plans come from the database misjudging how many rows "
    "will come back. So I looked at every case where it chose badly and "
    "measured the estimate. The median error was one point zero one seven. A "
    "perfect estimate is one point zero. In ninety eight percent of the bad "
    "choices the database knew the right numbers and still chose the slower "
    "path. That means fixing the estimates, which is what most tuning advice "
    "tries to do, cannot fix this."))
eyebrow(s, "Finding 2  /  RQ2")
heading(s, "The textbook explanation does not apply here")
bignum(s, "98.2%", "of bad choices had an essentially\ncorrect row estimate",
       L, Inches(2.2), Inches(4.4), size=72, color=RED)
grid(s, L, Inches(4.6),
     ["", "Value"],
     [["Bad choices examined", "113"],
      ["Median q-error among them", "1.017"],
      ["A perfect estimate", "1.000"]],
     [Inches(3.3), Inches(1.1)],
     cell_size=13.5, row_h=Inches(0.4))
figure(s, "fig3_qerror_vs_regret.png", Inches(7.5), Inches(1.95), Inches(4.9), Inches(4.5))
text(s, Inches(5.05), Inches(6.3), Inches(7.6), Inches(0.9),
     "The database knew the right numbers and chose badly anyway, so better "
     "statistics cannot be the remedy.",
     size=14, color=MUTED, spacing=1.3)
pagenum(s, 8)

# ================================================================ 9. F3 intro
s = slide(dark=True, notes=(
    "Given that, I tested the advice itself. Four widely recommended practices. "
    "Each one is sensible, each one appears in official documentation or "
    "well-regarded blogs. Under conditions I can identify, all four make the "
    "workload slower. I will take them one at a time."))
eyebrow(s, "Finding 3  /  RQ3")
heading(s, "Four recommended practices, measured", color=ONDARK)
line(s, L, Inches(2.1), CW, color=RGBColor(0x33, 0x38, 0x3D))
items = [
    ("Extended statistics", "for correlated columns", "Fixes the estimate 108x. Makes the query 2.7x slower."),
    ("Lower random_page_cost", "because the disk is an SSD", "The default is already within 6% of optimal. Lowering it is 2.3x worse."),
    ("Add a BRIN index", "it is small and cheap", "Bad-choice rate goes from 2.3% to 73.3%. Worst case 194x."),
    ("Enable Multi-Range Read", "MariaDB, ships switched off", "At the shipped buffer size the workload gets 7% slower."),
]
y = Inches(2.5)
for name, why, result in items:
    text(s, L, y, Inches(4.3), Inches(0.8), [
        [(name, {"size": 16, "bold": True, "color": ONDARK})],
        [(why, {"size": 12.5, "color": DIMDK})],
    ], spacing=1.3)
    text(s, Inches(6.0), y + Inches(0.06), Inches(6.4), Inches(0.6), result,
         size=15, color=RED_D, spacing=1.25)
    y += Inches(1.02)
    if y < Inches(6.4):
        line(s, L, y - Inches(0.16), CW, color=RGBColor(0x28, 0x2D, 0x32))
pagenum(s, 9)

# ================================================================ 10. F3a
s = slide(notes=(
    "First one. When two columns are correlated PostgreSQL misjudges them, and "
    "the official remedy is to create extended statistics. It works "
    "beautifully on the estimate. At full correlation it improves the estimate "
    "by a factor of one hundred and eight. And the query gets slower. Every "
    "single correlation level gets slower. The better estimate pushes the "
    "planner into a plan that costs less on paper and more in reality."))
eyebrow(s, "Finding 3a")
heading(s, "Extended statistics fix the estimate and damage the plan")
grid(s, L, Inches(2.1),
     ["Correlation", "q-error without", "with", "Estimate", "Time without", "with", "Result"],
     [["0.25", "25.2", "1.11", "23x better", "0.93 ms", "2.54 ms", "2.7x slower"],
      ["0.50", "49.5", "1.04", "48x better", "2.12 ms", "3.53 ms", "1.7x slower"],
      ["0.75", "75.9", "1.02", "75x better", "2.80 ms", "4.34 ms", "1.6x slower"],
      ["1.00", "113.2", "1.05", "108x better", "3.95 ms", "5.06 ms", "1.3x slower"]],
     [Inches(1.5), Inches(1.85), Inches(0.95), Inches(1.65), Inches(1.65), Inches(1.25), Inches(2.05)],
     aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 6,
     cell_size=13.5, row_h=Inches(0.40))
figure(s, "fig4_extended_statistics.png", Inches(2.4), Inches(4.15), Inches(8.5), Inches(2.45))
caption(s, "The remedy does exactly what it claims to the estimate. The plan still gets worse at every level.")
pagenum(s, 10)

# ================================================================ 11. F3b
s = slide(notes=(
    "Second one. Lowering random page cost because SSDs handle random reads "
    "well. I want to be careful here, because on one individual query I first "
    "measured a twenty four times effect and that was misleading. Measured "
    "across the whole workload, which is the honest way to do it, the shipped "
    "default of four is within six percent of the best setting, and lowering "
    "it to one makes the workload more than twice as slow."))
eyebrow(s, "Finding 3b")
heading(s, "The default random_page_cost is already near optimal")
grid(s, L, Inches(2.1),
     ["Setting", "Total workload time", "Versus best"],
     [["1.0", "2,306 ms", "2.3x worse"],
      ["1.1", "1,954 ms", "1.9x worse"],
      ["1.2", "1,942 ms", "1.9x worse"],
      ["1.5", "1,006 ms", "best"],
      ["2.0", "1,053 ms", "equal"],
      ["3.0", "1,053 ms", "equal"],
      ["4.0  (shipped default)", "1,068 ms", "within 6% of best"]],
     [Inches(2.45), Inches(1.95), Inches(1.75)],
     row_h=Inches(0.4), cell_size=13, hilite=[0, 6])
figure(s, "fig6_random_page_cost.png", Inches(7.35), Inches(2.25), Inches(5.05), Inches(2.9))
text(s, Inches(7.35), Inches(5.35), Inches(5.05), Inches(1.2), [
    [("A caution about how this is measured", {"size": 13, "bold": True, "color": RED})],
    [("On a single query I first measured a 24x effect. Across the workload it "
      "vanishes. Tuning verdicts drawn from one query do not survive.",
      {"size": 13, "color": MUTED})],
], spacing=1.3)
pagenum(s, 11)

# ================================================================ 12. F3c
s = slide(dark=True, notes=(
    "Third one, and this is the largest effect in the study. Adding a BRIN "
    "index alongside an existing B-tree, on the same column, because BRIN is "
    "tiny and people treat it as free. The bad-choice rate goes from two point "
    "three percent to seventy three point three percent. On one controlled "
    "test where the only thing I changed was whether the BRIN index existed, "
    "the query went from a third of a millisecond to sixty six milliseconds."))
eyebrow(s, "Finding 3c")
heading(s, "Adding a BRIN index because it is small and cheap", color=ONDARK)
bignum(s, "194x", "slower on a controlled test where the only change\nwas that a BRIN index existed alongside the B-tree",
       L, Inches(2.4), Inches(6.1), size=96, color=RED_D, lab_color=DIMDK)
text(s, L, Inches(5.05), Inches(6.1), Inches(0.6),
     "0.34 ms  becomes  65.95 ms", size=22, font=SERIF, color=ONDARK, spacing=1.2)
grid(s, Inches(7.3), Inches(2.5),
     ["Table size", "With BRIN", "Without", "Worst case"],
     [["1,000,000", "73.3%", "2.3%", "18.8x"],
      ["10,000,000", "30.1%", "32.3%", "2.6x"]],
     [Inches(1.9), Inches(1.25), Inches(1.05), Inches(1.15)],
     aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 3,
     cell_size=14, row_h=Inches(0.5))
text(s, Inches(7.4), Inches(4.35), Inches(4.9), Inches(2.2), [
    [("The bad-choice rate rises from 2.3% to 73.3% purely because a second "
      "index exists that is never the faster one to use.",
      {"size": 15, "color": ONDARK})],
    [("", {"size": 8})],
    [("The effect is confined to tables that still fit in the buffer pool, "
      "which is exactly the size at which BRIN is most often recommended as "
      "a free addition.", {"size": 14, "color": DIMDK})],
], spacing=1.35)
pagenum(s, 12)

# ================================================================ 13. CAUSE
s = slide(notes=(
    "I did not want to leave that as an observation, so I went into the "
    "PostgreSQL source. The function is choose_bitmap_and in indxpath.c. When "
    "two indexes can answer the same query it keeps one and throws the other "
    "away, and it decides by which index is cheaper to scan. But that cost is "
    "the cost of reading the index alone. It ignores the work the query then "
    "does on the table. A BRIN index is tiny by design, so it wins that "
    "comparison essentially always, and the index being discarded is the one "
    "that would have been almost two hundred times faster."))
eyebrow(s, "Finding 4  /  the mechanism")
heading(s, "Located in the source, not inferred")
text(s, L, Inches(2.05), Inches(6.3), Inches(0.5),
     "choose_bitmap_and()   in   optimizer/path/indxpath.c",
     size=15, font=MONO, color=BLUE, spacing=1.2)
text(s, L, Inches(2.7), Inches(6.3), Inches(2.6), [
    [("When two indexes can answer the same predicate, PostgreSQL keeps one "
      "and discards the other, choosing whichever is cheaper to scan.",
      {"size": 15.5})],
    [("", {"size": 7})],
    [("That cost is indextotalcost, the cost of reading the index alone. It "
      "does not include the heap fetches the query must then perform.",
      {"size": 15.5, "color": MUTED})],
    [("", {"size": 7})],
    [("A BRIN index is small by design, so it wins that comparison almost "
      "always, and the index thrown away is the faster one.",
      {"size": 15.5, "bold": True})],
], spacing=1.38)
grid(s, Inches(7.4), Inches(2.6),
     ["Index", "Cost used to decide", "Actual query time"],
     [["BRIN", "12.13   (kept)", "65.95 ms"],
      ["B-tree", "213.35   (discarded)", "0.34 ms"]],
     [Inches(1.3), Inches(2.3), Inches(1.75)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
     cell_size=14, row_h=Inches(0.52))
text(s, Inches(7.4), Inches(4.5), Inches(5.0), Inches(1.0),
     "The index that costs 17x less to read is the one that runs 194x slower.",
     size=15, color=RED, bold=True, spacing=1.3)
caption(s, "This is why the effect is a design consequence rather than a bug, and why it cannot be tuned away.")
pagenum(s, 13)

# ================================================================ 14. BOUNDARY
s = slide(notes=(
    "This is the slide I would most want to be asked about. The problem only "
    "happens while the table fits in the buffer pool, which you would expect. "
    "But look at the two columns separately. How often it goes wrong collapses "
    "right at the memory boundary, from seventy one percent down to thirty "
    "eight. How badly it goes wrong in the worst case does the opposite, it "
    "doubles at that same boundary and stays high all the way to three million "
    "rows. If I had tested only the smallest and the largest size, which is "
    "what most studies do, I would have concluded the problem simply "
    "disappears as tables grow. It does not."))
eyebrow(s, "Finding 5  /  RQ4")
heading(s, "The boundary is real, and it moves both ways at once")
grid(s, L, Inches(1.92),
     ["Table size", "Heap", "In pool", "Chose badly", "Worst case"],
     [["1,000,000", "112 MB", "yes", "71.4%", "18.8x"],
      ["1,250,000", "140 MB", "no", "37.9%", "38.8x"],
      ["1,500,000", "168 MB", "no", "31.0%", "40.3x"],
      ["2,000,000", "223 MB", "no", "25.8%", "38.3x"],
      ["3,000,000", "335 MB", "no", "30.3%", "41.5x"],
      ["5,000,000", "558 MB", "no", "36.4%", "7.0x"],
      ["10,000,000", "1,116 MB", "no", "38.2%", "1.9x"]],
     [Inches(1.75), Inches(1.15), Inches(1.05), Inches(1.2), Inches(1.15)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
     row_h=Inches(0.4), cell_size=13, header_h=Inches(0.56), hilite=[0, 4])
figure(s, "fig7_scale_transition.png", Inches(7.35), Inches(2.05), Inches(5.05), Inches(2.5))
text(s, Inches(7.35), Inches(4.75), Inches(5.05), Inches(1.9), [
    [("Two things happen at two different points",
      {"size": 13.5, "bold": True, "color": RED})],
    [("", {"size": 5})],
    [("How often it fails collapses at the memory boundary. How badly it fails "
      "doubles there and only subsides much later. Testing two sizes would "
      "have given the opposite conclusion.",
      {"size": 13.5, "color": MUTED})],
], spacing=1.3)
pagenum(s, 14)

# ================================================================ 15. CROSS
s = slide(notes=(
    "Adding the second database is what let me separate a PostgreSQL quirk "
    "from something general. The correlated-predicate failure turned out to be "
    "general. MariaDB estimates those predicates perfectly, exactly one point "
    "zero zero, right up until no single index covers both columns. Then it "
    "fails the same way PostgreSQL does. Two independently written databases "
    "failing identically means this is a design problem in how selectivity is "
    "composed, not a bug in either product."))
eyebrow(s, "Finding 6  /  RQ5")
heading(s, "What is general, and what belongs to one product")
grid(s, L, Inches(2.05),
     ["Correlation", "PostgreSQL q-error", "MariaDB q-error", "MariaDB plan used"],
     [["0.25", "26.4", "1.00", "single combined index"],
      ["0.50", "49.0", "1.00", "single combined index"],
      ["0.75", "81.8", "1.00", "single combined index"],
      ["1.00", "98.7", "55.8", "merges two indexes"]],
     [Inches(1.45), Inches(1.95), Inches(1.8), Inches(2.75)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.LEFT],
     cell_size=13.5, row_h=Inches(0.44), hilite=[3])
line(s, L, Inches(4.35), Inches(7.95))
text(s, L, Inches(4.65), Inches(7.95), Inches(2.0), [
    [("General:", {"size": 15.5, "bold": True, "color": RED}),
     ("  both engines are accurate only while a single index covers both "
      "columns and they can measure the answer directly. Once no such index "
      "exists, both multiply two separate guesses and both get it wrong.",
      {"size": 15.5})],
    [("", {"size": 7})],
    [("Product-specific:", {"size": 15.5, "bold": True, "color": BLUE}),
     ("  good plan selection. PostgreSQL's 2.3% bad-choice rate is not a "
      "property of relational optimisers in general.", {"size": 15.5})],
], spacing=1.38)
figure(s, "fig10_cross_conj.png", Inches(9.2), Inches(2.35), Inches(3.2), Inches(3.6))
pagenum(s, 15)

# ================================================================ 16. DEGRADE
s = slide(notes=(
    "And they degrade in opposite directions, which I think is the most useful "
    "practical result. As tables grow PostgreSQL gets wrong more often, up to "
    "thirty six percent, but its worst case never goes past three point seven "
    "times. MariaDB goes the other way, wrong less often but its worst case "
    "climbs to thirty one times. The reason is the plan repertoire. Bitmap "
    "scans put a ceiling on how bad a mistake can be because even a wrong "
    "choice still reads the disk in order. MariaDB has no ceiling. A database "
    "wrong often but never badly is easier to run than one wrong rarely but "
    "catastrophically, and one summary number hides that completely."))
eyebrow(s, "Finding 7")
heading(s, "The two engines degrade in opposite directions")
grid(s, L, Inches(2.0),
     ["Table size", "PostgreSQL wrong", "worst case", "MariaDB wrong", "worst case"],
     [["1,000,000", "0.0%", "1.1x", "52.9%", "3.3x"],
      ["1,250,000", "27.6%", "1.7x", "74.3%", "6.5x"],
      ["1,500,000", "6.9%", "1.8x", "77.8%", "10.9x"],
      ["2,000,000", "6.5%", "1.4x", "75.7%", "15.2x"],
      ["3,000,000", "6.1%", "1.7x", "36.1%", "18.8x"],
      ["5,000,000", "12.1%", "1.8x", "37.9%", "27.8x"],
      ["10,000,000", "36.1%", "3.7x", "12.4%", "30.9x"]],
     [Inches(1.6), Inches(1.3), Inches(1.1), Inches(1.3), Inches(1.15)],
     aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 4,
     row_h=Inches(0.4), cell_size=13, header_h=Inches(0.56), hilite=[6])
figure(s, "fig8_degradation_shape.png", Inches(7.6), Inches(2.05), Inches(4.8), Inches(2.4))
text(s, Inches(7.6), Inches(4.7), Inches(4.8), Inches(2.0), [
    [("Wrong often but never badly is easier to operate than wrong rarely but "
      "catastrophically.", {"size": 14, "bold": True})],
    [("", {"size": 6})],
    [("Bitmap scans cap the damage: even a wrong choice reads the disk in "
      "order. MariaDB has no such ceiling.", {"size": 13.5, "color": MUTED})],
], spacing=1.32)
caption(s, "Both engines judged only against plans the other could also have formed, so neither is held to a stricter standard.")
pagenum(s, 16)

# ================================================================ 17. COST
s = slide(notes=(
    "The computational cost, at ten million rows, counting only index types "
    "both databases have so the comparison is fair. PostgreSQL splits its time "
    "roughly evenly between building indexes and running queries. MariaDB "
    "builds in twenty minutes and then spends thirteen hours querying. The "
    "point I want to make is the last line: a study that measured only build "
    "time and a study that measured only query time would rank these two "
    "databases in opposite orders."))
eyebrow(s, "Finding 8")
heading(s, "Which cost dominates is a property of the engine")
grid(s, L, Inches(2.15),
     ["", "Time running queries", "Time building indexes"],
     [["PostgreSQL 17.1", "50.8 min", "45.6 min"],
      ["MariaDB 10.4.28", "781.8 min", "19.3 min"],
      ["Ratio", "15.4x slower", "2.4x faster"]],
     [Inches(2.5), Inches(2.15), Inches(2.15)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
     cell_size=14.5, row_h=Inches(0.5), hilite=[2])
line(s, L, Inches(4.25), Inches(6.9))
text(s, L, Inches(4.6), Inches(6.9), Inches(1.8), [
    [("A study measuring only build time and a study measuring only query time "
      "would rank these two databases in opposite orders.",
      {"size": 17, "font": SERIF})],
    [("", {"size": 8})],
    [("Measured at 10,000,000 rows with the same 128 MB pool, counting only "
      "index types both engines support.", {"size": 13.5, "color": MUTED})],
], spacing=1.35)
figure(s, "fig9_resource_cost.png", Inches(8.15), Inches(2.2), Inches(4.25), Inches(3.0))
pagenum(s, 17)

# ================================================================ 18. SELF-CORRECT
s = slide(notes=(
    "I want to show one thing I got wrong, because the correction is the "
    "better science. I had concluded that the correlated-column failure was "
    "caused by the plan shape that merges two indexes, since both databases "
    "fail while using it. Then I tested my own claim by switching that plan "
    "off entirely, which forces a completely different strategy. The error "
    "stayed the same. Fifty five point five nine against fifty five point "
    "eight four. So the merge is where the error becomes visible, not where it "
    "comes from. The real cause is narrower and more general, and I only found "
    "it by trying to refute myself."))
eyebrow(s, "Finding 10  /  validity")
heading(s, "One explanation I had to withdraw")
text(s, L, Inches(2.05), Inches(6.1), Inches(2.6), [
    [("What I claimed", {"size": 13, "bold": True, "color": MUTED})],
    [("That the correlated-column failure was caused by the plan shape that "
      "merges two indexes, since both engines fail while using it.",
      {"size": 15.5})],
    [("", {"size": 9})],
    [("How I tested it", {"size": 13, "bold": True, "color": MUTED})],
    [("Switched that plan off entirely, forcing a different strategy, and "
      "re-measured the estimation error.", {"size": 15.5})],
], spacing=1.36)
grid(s, Inches(7.2), Inches(2.15),
     ["Configuration", "q-error"],
     [["Merge plan enabled", "55.59"],
      ["Merge plan disabled", "55.84"]],
     [Inches(3.0), Inches(1.6)],
     cell_size=14.5, row_h=Inches(0.5))
text(s, Inches(7.2), Inches(4.0), Inches(5.2), Inches(1.2),
     "Unchanged. So the merge is where the error becomes visible, not where it comes from.",
     size=15, color=RED, bold=True, spacing=1.3)
line(s, L, Inches(5.15), CW)
text(s, L, Inches(5.45), CW, Inches(1.2), [
    [("The real cause is narrower: both engines are accurate only when one "
      "index covers both columns. Otherwise both multiply two separate "
      "guesses, whichever plan they choose.", {"size": 15.5})],
    [("", {"size": 6})],
    [("Three of my own measurement errors were found during the study. All "
      "three are reported in the paper rather than removed.",
      {"size": 14, "color": MUTED})],
], spacing=1.35)
pagenum(s, 18)

# ================================================================ 19. QC
s = slide(notes=(
    "Briefly on trustworthiness, because this is the part a marker will press "
    "on. Nineteen automated checks run before any measurement is accepted, "
    "including validating my generated data against PostgreSQL's own internal "
    "statistics, so the generator is checked against something I do not "
    "control. Everything is seeded, every measurement is written as it is "
    "taken so a power cut costs nothing, and the whole thing re-runs with one "
    "command."))
eyebrow(s, "Validity and reproducibility")
heading(s, "What makes these numbers trustworthy")
line(s, L, Inches(2.05), CW)
cols = [
    ("19", "automated checks run before any\nmeasurement is accepted",
     "Includes validating the generated data against PostgreSQL's own internal "
     "statistics, so the generator is checked against something I do not control."),
    ("7", "executions per query, first\ndiscarded as warm-up",
     "Parallel workers and JIT disabled, buffer pool pinned at 128 MB, dedicated "
     "tablespace, every dataset seeded so it regenerates identically."),
    ("3", "of my own errors found,\ncorrected and reported",
     "A collapsed selectivity range on skewed data, row counts meaning different "
     "things across the two engines, and an estimate that looked too good to be real."),
]
x = L
for num, lab, body in cols:
    text(s, x, Inches(2.45), Inches(3.5), Inches(1.0), num, size=58, font=SERIF,
         color=BLUE, spacing=0.95)
    text(s, x, Inches(3.35), Inches(3.5), Inches(0.7), lab, size=13, bold=True, spacing=1.25)
    text(s, x, Inches(4.25), Inches(3.5), Inches(2.0), body, size=13.5,
         color=MUTED, spacing=1.35)
    x += Inches(3.95)
caption(s, "Public repository with raw measurements, one-command reproduction, and a notebook that executes end to end.")
pagenum(s, 19)

# ================================================================ 20. SO WHAT
s = slide(notes=(
    "What it adds up to. Four pieces of common advice fail under conditions I "
    "can now name. The cause is not bad statistics, which is what almost all "
    "tuning targets, so the remedy has to be a costing change rather than a "
    "statistics change. And the practical rule I would give: never judge a "
    "tuning setting from one query, and never judge an optimiser from one "
    "table size."))
eyebrow(s, "What it means")
heading(s, "Conclusions")
line(s, L, Inches(2.05), CW)
concl = [
    ("Tuning advice needs conditions attached to it",
     "All four practices tested are sound in the case they were written for and harmful outside it. None of them ships with the boundary stated."),
    ("The cause is costing, not statistics",
     "98.2% of bad choices had correct row estimates. Better statistics cannot fix a comparison that ignores heap work, so the fix belongs in the cost model."),
    ("One number cannot describe an optimiser",
     "How often it errs and how badly it errs move in opposite directions across scale, and across these two engines."),
    ("Never generalise from one query or one table size",
     "A 24x tuning effect vanished across the workload, and a two-point scale test would have reversed the BRIN conclusion."),
]
y = Inches(2.4)
for i, (h, b) in enumerate(concl, 1):
    text(s, L, y, Inches(0.5), Inches(0.4), str(i), size=15, font=SERIF, bold=True,
         color=RED, spacing=1.0)
    text(s, L + Inches(0.55), y - Inches(0.02), Inches(11.0), Inches(1.0), [
        [(h, {"size": 16.5, "bold": True})],
        [(b, {"size": 14, "color": MUTED})],
    ], spacing=1.3)
    y += Inches(1.12)
pagenum(s, 20)

# ================================================================ 21. CLOSE
s = slide(dark=True, notes=(
    "That is the work. The code, the raw measurements and the notebook are all "
    "in the repository. Happy to take questions, and if you want to see any "
    "number in the deck verified I can show it coming out of the raw data."))
line(s, L, Inches(2.0), Inches(1.4), color=RED, weight=2.5)
text(s, L, Inches(2.35), Inches(9.5), Inches(1.4),
     "Thank you. Questions?", size=44, font=SERIF, color=ONDARK, spacing=1.1)
line(s, L, Inches(3.75), CW, color=RGBColor(0x33, 0x38, 0x3D))
text(s, L, Inches(4.1), Inches(5.6), Inches(2.2), [
    [("Everything is reproducible", {"size": 14, "bold": True, "color": ONDARK})],
    [("", {"size": 6})],
    [("23,480 raw measurements, seeded generators, the analysis notebook and "
      "both write-ups are in the repository.", {"size": 14, "color": DIMDK})],
], spacing=1.35)
text(s, Inches(7.2), Inches(4.1), Inches(5.2), Inches(2.2), [
    [("Md. Imtiaj Alam Sajin", {"size": 15, "bold": True, "color": ONDARK})],
    [("26-94090-2", {"size": 13, "color": DIMDK})],
    [("Supervisor: Dr. Ashraf Uddin", {"size": 13, "color": DIMDK})],
], spacing=1.4)
text(s, L, H - Inches(0.85), CW, Inches(0.4),
     "github.com/Imtiaj-Sajin/Reseach-on-Database-Architecture",
     size=11.5, font=MONO, color=DIMDK, spacing=1.2)

prs.save(OUT)
print("wrote", OUT)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
