"""
Build the paper presentation deck.

Same visual system as build_presentation.py, but restructured as a research
talk centred on the two-engine comparison: what the study finds that is
general to cost-based optimisers, and what belongs to one product.

Every number is copied from research/results, research/compare_dbms.py output
and paper/paper.tex. Nothing here is invented.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = r"g:\codes\Ass\DTMS"
FIG = os.path.join(ROOT, "research", "results", "figures")
OUT = os.path.join(ROOT, "submission", "DTMS_Paper_Presentation.pptx")

# ---------------------------------------------------------------- palette
INK    = RGBColor(0x1A, 0x1D, 0x21)
PAPER  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x70, 0x75, 0x7A)
FAINT  = RGBColor(0xA8, 0xAC, 0xB0)
RULE   = RGBColor(0xD9, 0xD5, 0xCF)
RED    = RGBColor(0xA8, 0x32, 0x2A)
RED_D  = RGBColor(0xE2, 0x76, 0x60)
BLUE   = RGBColor(0x1F, 0x4E, 0x5F)
DARK   = RGBColor(0x14, 0x17, 0x1A)
ONDARK = RGBColor(0xF2, 0xEF, 0xEA)
DIMDK  = RGBColor(0x92, 0x98, 0x9E)

SERIF, SANS, MONO = "Georgia", "Segoe UI", "Consolas"

W, H = Inches(13.333), Inches(7.5)
L = Inches(0.9)
R = W - Inches(0.9)
CW = R - L

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide(dark=False, notes=""):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK if dark else PAPER
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    s._dark = dark
    return s


def text(s, x, y, w, h, runs, size=16, font=SANS, color=None, bold=False,
         align=PP_ALIGN.LEFT, spacing=1.25, anchor=MSO_ANCHOR.TOP):
    if color is None:
        color = ONDARK if getattr(s, "_dark", False) else INK
    tf = s.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        paras = [[(runs, {})]]
    elif runs and isinstance(runs[0], tuple):
        paras = [runs]
    else:
        paras = runs
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        for t, ov in para:
            r = p.add_run()
            r.text = t
            f = r.font
            f.name = ov.get("font", font)
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.color.rgb = ov.get("color", color)


def line(s, x, y, w, color=None, weight=0.75):
    if color is None:
        color = RULE if not getattr(s, "_dark", False) else RGBColor(0x33, 0x38, 0x3D)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(weight))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False


def eyebrow(s, label, color=None):
    if color is None:
        color = RED if not getattr(s, "_dark", False) else DIMDK
    text(s, L, Inches(0.62), CW, Inches(0.3), label.upper(), size=10.5,
         bold=True, color=color, spacing=1.0)


def heading(s, title, y=Inches(1.0), size=29, color=None):
    text(s, L, y, CW, Inches(1.1), title, size=size, font=SERIF, color=color,
         spacing=1.12)


def caption(s, txt, y=None):
    text(s, L, y or (H - Inches(0.82)), CW, Inches(0.4), txt, size=10.5,
         color=MUTED if not getattr(s, "_dark", False) else DIMDK, spacing=1.2)


def grid(s, x, y, cols, rows, widths, aligns=None, header_size=11,
         cell_size=13.5, row_h=Inches(0.42), hilite=None, hi_color=None,
         header_h=Inches(0.33)):
    if aligns is None:
        aligns = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * (len(cols) - 1)
    if hi_color is None:
        hi_color = RED
    dark = getattr(s, "_dark", False)
    base = ONDARK if dark else INK
    lab = DIMDK if dark else MUTED
    PAD = Inches(0.10)
    xs, acc = [], x
    for wd in widths:
        xs.append(acc)
        acc += wd
    ws = [wd - 2 * PAD for wd in widths]
    xs = [cx + PAD for cx in xs]
    for cx, cwid, ctxt, al in zip(xs, ws, cols, aligns):
        text(s, cx, y, cwid, header_h, ctxt, size=header_size, bold=True,
             color=lab, align=al, spacing=1.1, anchor=MSO_ANCHOR.BOTTOM)
    ry = y + header_h
    line(s, x, ry, sum(widths),
         color=(base if not dark else RGBColor(0x4A, 0x50, 0x56)), weight=1.0)
    ry += Inches(0.1)
    for i, row in enumerate(rows):
        strong = hilite is not None and i in hilite
        for cx, cwid, ctxt, al in zip(xs, ws, row, aligns):
            text(s, cx, ry + Inches(0.06), cwid, Inches(0.34), str(ctxt),
                 size=cell_size, bold=strong,
                 color=(hi_color if strong else base), align=al, spacing=1.0)
        ry += row_h
        if i < len(rows) - 1:
            line(s, x, ry - Inches(0.04), sum(widths),
                 color=(RULE if not dark else RGBColor(0x2B, 0x30, 0x35)),
                 weight=0.6)
    line(s, x, ry - Inches(0.02), sum(widths),
         color=(base if not dark else RGBColor(0x4A, 0x50, 0x56)), weight=1.0)
    return ry


def figure(s, name, x, y, max_w, max_h):
    from PIL import Image
    path = os.path.join(FIG, name)
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = max_w, int(max_w / ar)
    if h > max_h:
        h, w = max_h, int(max_h * ar)
    s.shapes.add_picture(path, x + int((max_w - w) / 2), y, w, h)


def bignum(s, value, label, x, y, w, size=76, color=None, lab_color=None):
    dark = getattr(s, "_dark", False)
    if color is None:
        color = ONDARK if dark else INK
    if lab_color is None:
        lab_color = DIMDK if dark else MUTED
    text(s, x, y, w, Inches(1.4), value, size=size, font=SERIF, color=color,
         spacing=0.95)
    text(s, x, y + Inches(size / 72.0 * 1.02) + Inches(0.13), w, Inches(0.8),
         label, size=12.5, color=lab_color, spacing=1.25)


def pagenum(s, n):
    text(s, R - Inches(0.7), H - Inches(0.58), Inches(0.7), Inches(0.3), str(n),
         size=10, align=PP_ALIGN.RIGHT, spacing=1.0,
         color=FAINT if not getattr(s, "_dark", False)
         else RGBColor(0x4A, 0x50, 0x56))


# ================================================================ 1 TITLE
s = slide(dark=True, notes=(
    "This is an empirical study of how two relational engines choose an "
    "access path, and whether the tuning advice written for that decision "
    "actually works. The reason there are two engines rather than one is the "
    "core of the talk: with a single system you cannot tell a property of "
    "cost-based optimisation from a quirk of one product."))
line(s, L, Inches(1.5), Inches(1.4), color=RED, weight=2.5)
text(s, L, Inches(1.8), Inches(10.8), Inches(2.3),
     "Measuring the Impact of Index Types and Tuning Settings on "
     "Query Execution Plans in PostgreSQL and MariaDB",
     size=35, font=SERIF, color=ONDARK, spacing=1.14)
text(s, L, Inches(4.3), Inches(10.8), Inches(0.4),
     "An access-path study across two independently written optimisers",
     size=14, color=DIMDK, spacing=1.2)
line(s, L, Inches(4.9), CW, color=RGBColor(0x33, 0x38, 0x3D))
text(s, L, Inches(5.2), Inches(5.4), Inches(1.3), [
    [("Md. Imtiaj Alam Sajin", {"size": 15, "bold": True, "color": ONDARK})],
    [("26-94090-2   |   26-94090-2@student.aiub.edu", {"size": 12, "color": DIMDK})],
], spacing=1.45)
text(s, Inches(7.2), Inches(5.2), Inches(5.2), Inches(1.3), [
    [("Supervisor", {"size": 11, "color": DIMDK})],
    [("Dr. Ashraf Uddin", {"size": 15, "bold": True, "color": ONDARK})],
], spacing=1.45)
text(s, L, H - Inches(0.85), CW, Inches(0.4),
     "23,480 measurements  |  PostgreSQL 17.1 and MariaDB 10.4.28  |  "
     "seven table scales  |  fully reproducible",
     size=11.5, font=MONO, color=DIMDK, spacing=1.2)

# ================================================================ 2 PREMISE
s = slide(notes=(
    "The starting point. Every query needs one decision made for it, and the "
    "same query on the same data can be a hundred times slower depending on "
    "it. Because the decision is invisible, a large body of tuning advice has "
    "grown around helping the database make it. That advice is written by "
    "good engineers and it sounds right. Almost none of it has been measured "
    "under controlled conditions, so nobody knows the conditions under which "
    "it holds."))
eyebrow(s, "Motivation")
heading(s, "A decision nobody sees, and advice nobody measured", size=31)
line(s, L, Inches(2.2), CW)
text(s, L, Inches(2.55), Inches(6.1), Inches(3.0), [
    [("Scan the table, or use an index, and if so which one.",
      {"size": 17, "bold": True})],
    [("", {"size": 8})],
    [("One decision per query, invisible to the user, worth up to two orders "
      "of magnitude in runtime.", {"size": 15.5, "color": MUTED})],
    [("", {"size": 8})],
    [("The literature attributes bad plans to cardinality misestimation. "
      "That work targets join ordering. The single-relation access-path "
      "decision has not been isolated the same way.",
      {"size": 15.5, "color": MUTED})],
], spacing=1.4)
text(s, Inches(7.6), Inches(2.55), Inches(4.8), Inches(3.4), [
    [("Four widely recommended practices", {"size": 13, "bold": True, "color": RED})],
    [("", {"size": 7})],
    [("Extended statistics for correlated columns", {"size": 14.5})],
    [("Lower random_page_cost on SSD storage", {"size": 14.5})],
    [("Add a BRIN index, it costs almost nothing", {"size": 14.5})],
    [("Enable Multi-Range Read in MariaDB", {"size": 14.5})],
    [("", {"size": 7})],
    [("Every one is sound in the case it was written for. "
      "This study measures where each stops being sound.",
      {"size": 14.5, "bold": True})],
], spacing=1.45)
pagenum(s, 2)

# ================================================================ 3 WHY TWO
s = slide(dark=True, notes=(
    "This is the slide that explains why the study has two databases rather "
    "than one, and it is the frame for everything after it. If you measure "
    "one optimiser and find a failure, you cannot tell whether you have found "
    "something about cost-based optimisation or a bug in one product. Two "
    "independently written engines let you separate those. When both fail the "
    "same way it points at a shared design assumption. When only one fails it "
    "is an implementation choice. Every result from here on is labelled one "
    "way or the other."))
eyebrow(s, "The central move")
heading(s, "Why two engines and not one", color=ONDARK)
line(s, L, Inches(2.1), CW, color=RGBColor(0x33, 0x38, 0x3D))
text(s, L, Inches(2.5), Inches(11.5), Inches(0.8),
     "A single-system study cannot separate a property of cost-based "
     "access-path selection from an artefact of one implementation.",
     size=20, font=SERIF, color=ONDARK, spacing=1.25)
line(s, L, Inches(3.75), CW, color=RGBColor(0x33, 0x38, 0x3D))
text(s, L, Inches(4.1), Inches(5.6), Inches(2.2), [
    [("Both engines fail the same way", {"size": 16, "bold": True, "color": RED_D})],
    [("", {"size": 6})],
    [("Points to a shared design assumption. PostgreSQL and MariaDB were "
      "written independently, so agreement is evidence, not coincidence.",
      {"size": 14.5, "color": DIMDK})],
], spacing=1.35)
text(s, Inches(7.2), Inches(4.1), Inches(5.2), Inches(2.2), [
    [("Only one fails", {"size": 16, "bold": True, "color": ONDARK})],
    [("", {"size": 6})],
    [("An implementation choice, and the other engine shows what a different "
      "choice would have bought.", {"size": 14.5, "color": DIMDK})],
], spacing=1.35)
caption(s, "Both systems run identical generated data, identical queries, identical metrics and a matched 128 MB buffer pool.")
pagenum(s, 3)

# ================================================================ 4 METRIC
s = slide(notes=(
    "Two definitions. Regret is how many times slower the plan the engine "
    "chose was than the fastest plan actually available to it, measured by "
    "forcing each index configuration separately and re-running the same "
    "query. A regret of one means it did the best it could. That separates "
    "the quality of the decision from the raw speed of the engine, which "
    "matters when comparing two products. Q-error is the standard measure of "
    "estimation accuracy. The whole paper turns on these two being reported "
    "separately."))
eyebrow(s, "Metrics")
heading(s, "Two measures, deliberately kept apart")
text(s, L, Inches(2.1), Inches(7.6), Inches(0.6),
     "R  =  time of the plan chosen  /  time of the fastest plan available",
     size=16, font=SERIF, color=BLUE, spacing=1.2)
line(s, L, Inches(2.85), Inches(7.6))
text(s, L, Inches(3.15), Inches(7.6), Inches(2.6), [
    [("Access-path regret", {"size": 14, "bold": True}),
     ("  measures the decision. The comparison is against plans the engine "
      "really had, forced one at a time, so a slow engine is not penalised "
      "for being slow.", {"size": 15})],
    [("", {"size": 8})],
    [("q-error", {"size": 14, "bold": True}),
     ("  measures the estimate: the factor by which the predicted row count "
      "misses the truth. Reported separately throughout, because the paper's "
      "main negative result is that the two come apart.", {"size": 15})],
], spacing=1.38)
bignum(s, "1.0", "the engine made the best\nchoice available to it",
       Inches(9.1), Inches(2.3), Inches(3.3), size=60, color=BLUE)
bignum(s, "23,480", "measurements across\nboth engines",
       Inches(9.1), Inches(4.5), Inches(3.3), size=52)
pagenum(s, 4)

# ================================================================ 5 SETUP
s = slide(notes=(
    "The design in one slide. Data is generated rather than taken from a "
    "benchmark, because a fixed dataset fixes skew, correlation and "
    "clustering at whatever values it happens to have, and then you cannot "
    "attribute an effect to any one property. Generating it also means the "
    "true row count is known for every query, so misestimation is measured "
    "rather than inferred. Both engines cover the same seven table sizes, so "
    "the comparison is symmetric rather than one system being probed harder."))
eyebrow(s, "Experimental design")
heading(s, "Symmetric, controlled, and seeded")
text(s, L, Inches(2.0), Inches(5.5), Inches(2.6), [
    [("Three properties varied independently", {"size": 14, "bold": True})],
    [("", {"size": 6})],
    [("Value skew (Zipfian, 4 levels), predicate correlation (5 levels), "
      "physical clustering (4 levels), giving 11 dataset configurations.",
      {"size": 14.5, "color": MUTED})],
    [("", {"size": 8})],
    [("Why not TPC-H", {"size": 14, "bold": True})],
    [("A fixed dataset fixes all three at whatever values it happens to "
      "have, so no effect can be attributed to a single property. Generating "
      "the data also makes the true row count known for every query.",
      {"size": 14.5, "color": MUTED})],
], spacing=1.36)
grid(s, Inches(6.9), Inches(2.0),
     ["Table size", "PostgreSQL", "MariaDB"],
     [["1,000,000", "2,398", "1,232"],
      ["1,250,000", "654", "336"],
      ["1,500,000", "654", "336"],
      ["2,000,000", "654", "336"],
      ["3,000,000", "654", "336"],
      ["5,000,000", "654", "336"],
      ["10,000,000", "2,398", "1,232"],
      ["Configuration sweeps", "6,230", "5,040"],
      ["Total", "14,296", "9,184"]],
     [Inches(2.7), Inches(1.5), Inches(1.4)],
     row_h=Inches(0.395), cell_size=13, hilite=[8], hi_color=INK)
caption(s, "10 index configurations, 4 query families, 8 selectivity targets. Seven runs per query, first discarded. Parallelism and JIT off.")
pagenum(s, 5)

# ================================================================ 6 RQ1
s = slide(notes=(
    "The first result and the frame for the comparison. On identical data "
    "PostgreSQL picks the best available plan 97.7 percent of the time. That "
    "matters because it means the PostgreSQL failures I show later are "
    "specific, identifiable defects rather than a generally weak optimiser. "
    "MariaDB errs on two thirds of queries. But notice the last column: its "
    "median error is mild. Wrong often, mostly by a little."))
eyebrow(s, "Result 1")
heading(s, "Selection quality differs by an order of magnitude")
grid(s, L, Inches(2.15),
     ["Engine", "Queries", "Chose badly", "Median regret", "Worst case"],
     [["PostgreSQL 17.1", "130", "2.3%", "1.01", "1.39"],
      ["MariaDB 10.4.28", "153", "66.0%", "1.34", "7.91"]],
     [Inches(2.9), Inches(1.35), Inches(1.6), Inches(1.7), Inches(1.5)],
     aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 4,
     row_h=Inches(0.48), cell_size=14.5)
line(s, L, Inches(3.85), Inches(9.05))
text(s, L, Inches(4.15), Inches(6.1), Inches(2.4), [
    [("Read the last two columns together", {"size": 14, "bold": True, "color": RED})],
    [("", {"size": 6})],
    [("PostgreSQL is almost always right, and when wrong is wrong by 39%. "
      "MariaDB is usually wrong, but its median error is 34%. These are two "
      "different failure profiles, not one engine being better at everything, "
      "and the rest of the talk separates them.",
      {"size": 14.5, "color": MUTED})],
], spacing=1.38)
figure(s, "fig2_regret_by_selectivity.png", Inches(7.3), Inches(4.0), Inches(5.1), Inches(2.6))
caption(s, "Measured at 1,000,000 rows. Both engines judged only against plans the other could also have formed.")
pagenum(s, 6)

# ================================================================ 7 PER FAMILY
s = slide(notes=(
    "Breaking that down by query family shows the difference is not uniform. "
    "PostgreSQL's only real weakness is the conjunctive family, ten percent, "
    "and note its estimation error there is 71, by far the worst number in "
    "its column. MariaDB is the mirror image. Its estimates on conjunctions "
    "are perfect, exactly 1.00, yet it still chooses badly 45 percent of the "
    "time. And look at equality: three queries, wrong on all three, median "
    "regret 6.8. Perfect knowledge, wrong decision."))
eyebrow(s, "Result 1, decomposed")
heading(s, "The two engines fail on different families")
grid(s, L, Inches(2.0),
     ["Family", "PG n", "PG wrong", "PG median", "PG q-error",
      "MDB n", "MDB wrong", "MDB median", "MDB q-error"],
     [["conj", "19", "10.5%", "1.044", "71.3", "20", "45.0%", "1.067", "1.00"],
      ["eq", "3", "0.0%", "1.013", "1.00", "3", "100.0%", "6.837", "1.29"],
      ["range", "55", "1.8%", "1.013", "1.01", "66", "71.2%", "1.347", "1.84"],
      ["ts_range", "53", "0.0%", "1.001", "1.01", "64", "65.6%", "1.343", "1.97"]],
     [Inches(1.35), Inches(0.75), Inches(1.2), Inches(1.25), Inches(1.25),
      Inches(0.85), Inches(1.3), Inches(1.35), Inches(1.35)],
     aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 8,
     cell_size=13, header_size=10, header_h=Inches(0.55), row_h=Inches(0.44),
     hilite=[1])
line(s, L, Inches(4.5), CW)
text(s, L, Inches(4.85), Inches(5.9), Inches(1.9), [
    [("PostgreSQL", {"size": 14, "bold": True, "color": BLUE}),
     ("  errs almost only on conjunctions, and that is exactly where its "
      "estimation error is worst, 71.3. Here the textbook story holds.",
      {"size": 14.5})],
], spacing=1.35)
text(s, Inches(7.2), Inches(4.85), Inches(5.2), Inches(1.9), [
    [("MariaDB", {"size": 14, "bold": True, "color": RED}),
     ("  estimates conjunctions perfectly at 1.00 and still errs on 45% of "
      "them. On equality it is wrong every time, by a median of 6.8x, with a "
      "near-correct estimate.", {"size": 14.5})],
], spacing=1.35)
pagenum(s, 7)

# ================================================================ 8 REPERTOIRE
s = slide(notes=(
    "Here is the mechanism behind that, and it is the most important "
    "structural result in the comparison. This counts what each engine "
    "actually chose when given every index. PostgreSQL answers almost "
    "everything with a bitmap scan: it collects the matching row locations, "
    "sorts them into physical order, then reads the heap in order. MariaDB "
    "never forms one except seven times on conjunctions. It either does a "
    "direct index lookup or gives up and scans the whole table, sixty six "
    "times, where PostgreSQL scanned nine. The difference in outcomes is a "
    "difference in the plans available, not in how carefully each engine "
    "costs them."))
eyebrow(s, "Mechanism")
heading(s, "The gap is a difference in plan repertoire")
grid(s, L, Inches(2.05),
     ["Family", "PG bitmap", "PG index", "PG seq",
      "MDB bitmap", "MDB index", "MDB seq"],
     [["conj", "110", "0", "0", "7", "103", "0"],
      ["eq", "88", "0", "0", "0", "88", "0"],
      ["range", "88", "0", "0", "0", "55", "33"],
      ["ts_range", "60", "19", "9", "0", "55", "33"]],
     [Inches(1.5), Inches(1.4), Inches(1.2), Inches(1.05),
      Inches(1.55), Inches(1.35), Inches(1.2)],
     aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 6,
     cell_size=13.5, header_size=10.5, header_h=Inches(0.55),
     row_h=Inches(0.44))
line(s, L, Inches(4.5), Inches(9.25))
text(s, L, Inches(4.85), Inches(6.0), Inches(2.0), [
    [("PostgreSQL sorts row locations into physical order before touching "
      "the heap, so even a wrong choice reads the disk sequentially. That "
      "puts a ceiling on how bad a mistake can be.", {"size": 15})],
    [("", {"size": 7})],
    [("MariaDB has no equivalent enabled. A mistaken index scan becomes "
      "random single-row fetches, and there is no ceiling.",
      {"size": 15, "bold": True})],
], spacing=1.35)
text(s, Inches(9.6), Inches(4.85), Inches(2.8), Inches(2.0), [
    [("On conjunctions MariaDB's", {"size": 13, "color": MUTED})],
    [("raw access types are", {"size": 13, "color": MUTED})],
    [("ref 103, index_merge 7.", {"size": 13, "font": MONO, "color": INK})],
    [("It merges two indexes only", {"size": 13, "color": MUTED})],
    [("at full correlation.", {"size": 13, "color": MUTED})],
], spacing=1.3)
pagenum(s, 8)

# ================================================================ 9 WORST ROW
s = slide(dark=True, notes=(
    "One row from the raw data makes the point concretely. This is MariaDB's "
    "single worst measurement. A query returning ten percent of a million "
    "rows. It estimated 218 thousand against a true 101 thousand, so the "
    "estimate was off by only a factor of two, which for an optimiser is "
    "essentially correct. It chose an index scan. Scanning the whole table "
    "would have taken 249 milliseconds. It took 1,970. The estimate was fine. "
    "The decision was not."))
eyebrow(s, "One row from the raw data")
heading(s, "MariaDB's worst case was not an estimation failure", color=ONDARK)
grid(s, L, Inches(2.3),
     ["Field", "Value"],
     [["dataset / family", "skew10 / eq"],
      ["true rows", "101,636  (10.2% of the table)"],
      ["estimated rows", "218,632"],
      ["q-error", "2.151   (near-correct for an optimiser)"],
      ["path chosen", "index scan"],
      ["best path available", "sequential scan"],
      ["time taken", "1,969.7 ms"],
      ["time if scanned", "248.9 ms"],
      ["regret", "7.92x"]],
     [Inches(2.6), Inches(4.5)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT],
     cell_size=13.5, row_h=Inches(0.42), hilite=[3, 8], hi_color=RED_D)
text(s, Inches(8.3), Inches(2.6), Inches(4.1), Inches(3.4), [
    [("It knew roughly how many rows were coming back and still chose to "
      "fetch a tenth of the table one row at a time.",
      {"size": 16, "color": ONDARK})],
    [("", {"size": 10})],
    [("Fixing the statistics would not have changed this decision. The cost "
      "model, not the estimate, is what ranked the index scan first.",
      {"size": 14.5, "color": DIMDK})],
], spacing=1.4)
pagenum(s, 9)

# ================================================================ 10 RQ3
s = slide(notes=(
    "That was not an isolated case. Across every bad choice in PostgreSQL, "
    "the median estimation error is 1.017, where a perfect estimate is 1.0. "
    "In 98 percent of the cases where the optimiser chose a slower path, it "
    "had the row count essentially right. This is the paper's main negative "
    "result and it cuts against the standard account, which attributes bad "
    "plans to misestimation. For this decision, better statistics cannot be "
    "the remedy, because the statistics were already correct."))
eyebrow(s, "Result 2, the main negative result")
heading(s, "Misestimation does not explain misselection")
bignum(s, "98.2%", "of bad access-path choices had an\nessentially correct row estimate",
       L, Inches(2.2), Inches(4.6), size=70, color=RED)
grid(s, L, Inches(4.75),
     ["", "Value"],
     [["Bad choices examined", "113"],
      ["Median q-error among them", "1.017"],
      ["A perfect estimate", "1.000"]],
     [Inches(3.4), Inches(1.2)],
     cell_size=13.5, row_h=Inches(0.4))
figure(s, "fig3_qerror_vs_regret.png", Inches(7.6), Inches(1.95), Inches(4.8), Inches(4.4))
text(s, Inches(5.2), Inches(6.35), Inches(7.2), Inches(0.9),
     "The remedy most tuning advice reaches for is aimed at a cause that is "
     "not operating here.", size=14.5, color=MUTED, spacing=1.3)
pagenum(s, 10)

# ================================================================ 11 GENERAL
s = slide(notes=(
    "Now the two-engine payoff, and this is the result I would most want a "
    "reviewer to look at. Correlated conjunctions. PostgreSQL's error climbs "
    "with correlation, from 1.13 up to 98.65. MariaDB is perfect, exactly "
    "1.00, all the way up to full correlation, where it jumps to 55.84. The "
    "right-hand columns are the control: the same queries on independent "
    "columns, where both stay near one until the very end. Two independently "
    "written engines, accurate while a single index covers both columns, both "
    "failing once it does not. That makes this a design problem, not a bug in "
    "either product."))
eyebrow(s, "Result 3  /  what is general")
heading(s, "Both engines fail once no single index covers both columns")
grid(s, L, Inches(2.05),
     ["Correlation", "PostgreSQL dependent", "independent",
      "MariaDB dependent", "independent"],
     [["0.00", "1.13", "1.08", "1.00", "1.00"],
      ["0.25", "26.39", "1.37", "1.00", "1.00"],
      ["0.50", "49.01", "2.14", "1.00", "1.00"],
      ["0.75", "81.82", "3.44", "1.00", "1.00"],
      ["1.00", "98.65", "98.00", "55.84", "1.00"]],
     [Inches(1.4), Inches(1.95), Inches(1.4), Inches(1.85), Inches(1.4)],
     aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 4,
     cell_size=13.5, header_size=10.5, header_h=Inches(0.55),
     row_h=Inches(0.42), hilite=[4])
line(s, L, Inches(4.75), Inches(9.1))
text(s, L, Inches(5.1), Inches(6.1), Inches(1.9), [
    [("The independent columns are the control.", {"size": 14.5, "bold": True}),
     ("  Same queries, uncorrelated data. Both engines stay near 1.0, so the "
      "error tracks correlation rather than the query shape.",
      {"size": 14.5})],
], spacing=1.35)
text(s, Inches(7.4), Inches(5.1), Inches(5.0), Inches(1.9), [
    [("Both are accurate only while one index measures both columns "
      "directly. Once none does, both multiply two separate guesses and both "
      "get it wrong.", {"size": 14.5, "bold": True, "color": RED})],
], spacing=1.35)
figure(s, "fig10_cross_conj.png", Inches(9.35), Inches(2.05), Inches(3.05), Inches(2.6))
pagenum(s, 11)

# ================================================================ 12 SPECIFIC
s = slide(dark=True, notes=(
    "So the study sorts its own findings into two piles, and this is what I "
    "mean by what the research truly identifies. On the left, things that "
    "showed up in both engines and are therefore properties of cost-based "
    "access-path selection. On the right, things only one engine does, which "
    "are implementation choices, and the other engine shows what the "
    "alternative buys you."))
eyebrow(s, "The classification")
heading(s, "What generalises, and what belongs to one product", color=ONDARK)
line(s, L, Inches(2.1), CW, color=RGBColor(0x33, 0x38, 0x3D))
text(s, L, Inches(2.5), Inches(5.7), Inches(0.4),
     "General to cost-based optimisers", size=15, bold=True, color=RED_D)
for i, (h, b) in enumerate([
    ("Per-column composition fails under correlation",
     "Both engines accurate to 1.00 with one covering index, both wrong without one."),
    ("Estimation quality and plan quality come apart",
     "Correct row counts accompany 98.2% of PostgreSQL's bad choices."),
    ("Advice ships without its boundary",
     "All four practices tested help in one regime and harm in another."),
]):
    y = Inches(3.0) + Inches(1.15) * i
    text(s, L, y, Inches(5.7), Inches(1.0), [
        [(h, {"size": 14.5, "bold": True, "color": ONDARK})],
        [(b, {"size": 13, "color": DIMDK})],
    ], spacing=1.3)
line(s, Inches(6.75), Inches(2.5), Pt(0.75), color=RGBColor(0x33, 0x38, 0x3D))
text(s, Inches(7.2), Inches(2.5), Inches(5.2), Inches(0.4),
     "Specific to one implementation", size=15, bold=True, color=ONDARK)
for i, (h, b) in enumerate([
    ("Good selection quality is PostgreSQL's, not the norm",
     "2.3% versus 66.0% on identical data. Not a property of relational optimisers."),
    ("The damage ceiling comes from bitmap scans",
     "MariaDB has no equivalent enabled, so its worst case is unbounded by scale."),
    ("The BRIN misselection is a PostgreSQL costing defect",
     "MariaDB cannot form the plan at all, so it cannot make this mistake."),
]):
    y = Inches(3.0) + Inches(1.15) * i
    text(s, Inches(7.2), y, Inches(5.2), Inches(1.0), [
        [(h, {"size": 14.5, "bold": True, "color": ONDARK})],
        [(b, {"size": 13, "color": DIMDK})],
    ], spacing=1.3)
pagenum(s, 12)

# ================================================================ 13 DEGRADE
s = slide(notes=(
    "Measuring both across all seven scales shows they do not merely differ "
    "by degree, they degrade in opposite directions. As tables grow "
    "PostgreSQL gets wrong more often, from never up to 36 percent, but its "
    "worst case never exceeds 3.7. MariaDB goes the other way: wrong less "
    "often, down from 53 to 12 percent, while its worst case climbs to 31. "
    "The operational reading matters. A database wrong often but never badly "
    "is easier to run than one wrong rarely but catastrophically, and "
    "summarising either with a single number hides exactly this."))
eyebrow(s, "Result 4")
heading(s, "The two engines degrade in opposite directions")
grid(s, L, Inches(2.0),
     ["Table size", "PostgreSQL wrong", "worst case", "MariaDB wrong", "worst case"],
     [["1,000,000", "0.0%", "1.1x", "52.9%", "3.3x"],
      ["1,250,000", "27.6%", "1.7x", "74.3%", "6.5x"],
      ["1,500,000", "6.9%", "1.8x", "77.8%", "10.9x"],
      ["2,000,000", "6.5%", "1.4x", "75.7%", "15.2x"],
      ["3,000,000", "6.1%", "1.7x", "36.1%", "18.8x"],
      ["5,000,000", "12.1%", "1.8x", "37.9%", "27.8x"],
      ["10,000,000", "36.1%", "3.7x", "12.4%", "30.9x"]],
     [Inches(1.6), Inches(1.3), Inches(1.1), Inches(1.3), Inches(1.15)],
     aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 4,
     row_h=Inches(0.4), cell_size=13, header_h=Inches(0.56), hilite=[6])
figure(s, "fig8_degradation_shape.png", Inches(7.6), Inches(2.05), Inches(4.8), Inches(2.4))
text(s, Inches(7.6), Inches(4.7), Inches(4.8), Inches(2.0), [
    [("Wrong often but never badly is easier to operate than wrong rarely "
      "but catastrophically.", {"size": 14.5, "bold": True})],
    [("", {"size": 6})],
    [("Frequency and severity are separate questions, and here they do not "
      "even move in the same direction.", {"size": 13.5, "color": MUTED})],
], spacing=1.32)
caption(s, "Both engines judged only against plans the other could also have formed, so neither is held to a stricter standard.")
pagenum(s, 13)

# ================================================================ 14 COST
s = slide(notes=(
    "The resource picture inverts too. At ten million rows, counting only "
    "index types both engines have, PostgreSQL splits its time roughly evenly "
    "between building indexes and running queries. MariaDB builds in twenty "
    "minutes and then spends thirteen hours querying. The consequence is on "
    "the bottom line: a study that measured only build time and a study that "
    "measured only query time would rank these two databases in opposite "
    "orders. Which cost dominates is a property of the engine, not of the "
    "workload."))
eyebrow(s, "Result 5")
heading(s, "Which cost dominates is a property of the engine")
grid(s, L, Inches(2.15),
     ["", "Time running queries", "Time building indexes"],
     [["PostgreSQL 17.1", "50.8 min", "45.6 min"],
      ["MariaDB 10.4.28", "781.8 min", "19.3 min"],
      ["Ratio", "15.4x slower", "2.4x faster"]],
     [Inches(2.5), Inches(2.15), Inches(2.15)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
     cell_size=14.5, row_h=Inches(0.5), hilite=[2])
line(s, L, Inches(4.25), Inches(6.9))
text(s, L, Inches(4.6), Inches(6.9), Inches(1.9), [
    [("A study measuring only build time and a study measuring only query "
      "time would rank these two engines in opposite orders.",
      {"size": 17, "font": SERIF})],
    [("", {"size": 8})],
    [("Measured at 10,000,000 rows, matched 128 MB pools, counting only index "
      "types both engines support.", {"size": 13.5, "color": MUTED})],
], spacing=1.35)
figure(s, "fig9_resource_cost.png", Inches(8.15), Inches(2.2), Inches(4.25), Inches(3.0))
pagenum(s, 14)

# ================================================================ 15 SWEEP
s = slide(notes=(
    "The obvious objection is that MariaDB was left on defaults. It does have "
    "a comparable mechanism, Multi-Range Read, which sorts row identifiers "
    "before fetching. It ships switched off. If turning it on closed the gap, "
    "the whole difference would be a configuration choice rather than a "
    "design one, so I swept all the settings that govern these decisions, "
    "five thousand and forty more measurements. Nothing closes the gap. "
    "Turning MRR on makes the workload seven percent slower, and I verified "
    "it genuinely engaged by checking the plan output. The last column is "
    "worth noting too: the estimation error sits at 55.84 in every "
    "configuration but one."))
eyebrow(s, "Result 6  /  ruling out the obvious objection")
heading(s, "The deficit cannot be configured away")
grid(s, L, Inches(2.05),
     ["MariaDB configuration", "Total time", "vs default", "Conjunction q-error"],
     [["ucs3  (statistics depth)", "100.7 s", "0.87", "55.84"],
      ["ucs5", "102.1 s", "0.88", "55.84"],
      ["no_index_merge", "107.0 s", "0.92", "55.59"],
      ["mrr_on_costbased", "108.4 s", "0.93", "55.84"],
      ["default", "116.1 s", "1.00", "55.84"],
      ["mrr_on_sorted", "121.9 s", "1.05", "55.84"],
      ["mrr_on", "123.8 s", "1.07", "55.84"]],
     [Inches(2.95), Inches(1.35), Inches(1.25), Inches(1.75)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
     row_h=Inches(0.4), cell_size=13, header_size=10.5, hilite=[4, 6])
text(s, Inches(8.5), Inches(2.15), Inches(3.9), Inches(4.2), [
    [("Enabling MRR makes the workload 7% slower.",
      {"size": 15, "bold": True, "color": RED})],
    [("", {"size": 8})],
    [("Verified that it genuinely engaged rather than being ignored: EXPLAIN "
      "reports a rowid-ordered scan and the plan JSON carries the marker, "
      "neither of which appears with MRR off.",
      {"size": 13.5, "color": MUTED})],
    [("", {"size": 8})],
    [("Statistics depth changes total time by at most 13% and leaves the "
      "conjunction error untouched. MariaDB's accuracy there comes from index "
      "dives, not from the statistical machinery.",
      {"size": 13.5, "color": MUTED})],
], spacing=1.32)
caption(s, "1,000,000 rows, five datasets, 5,040 measurements. No configuration closes the gap with PostgreSQL.")
pagenum(s, 15)

# ================================================================ 16 MRR BUF
s = slide(notes=(
    "And the reason MRR does not close the gap is a size limit, which is "
    "worth showing because it turns a vague claim into a bounded one. MRR "
    "sorts within a buffer that defaults to 256 kilobytes; a bitmap scan "
    "materialises the entire qualifying set. Raising the buffer to eight "
    "megabytes does help, but only on the narrowest queries: 0.65 at one "
    "percent selectivity, gone by five percent, nothing at ten. At its very "
    "best it recovers about a third of one query class against a cross-system "
    "gap of fifteen times. So this is a fourth piece of advice that backfires "
    "at its shipped setting."))
eyebrow(s, "Result 6, continued")
heading(s, "Why it cannot close it: the buffer is bounded")
grid(s, L, Inches(2.15),
     ["mrr_buffer_size", "1% selectivity", "5%", "10%"],
     [["MRR disabled", "1.00", "1.00", "1.00"],
      ["256 KB  (shipped default)", "1.44", "1.27", "1.02"],
      ["8 MB", "0.65", "1.48", "1.05"],
      ["64 MB", "0.66", "0.97", "1.00"]],
     [Inches(3.1), Inches(1.7), Inches(1.2), Inches(1.2)],
     aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 3,
     cell_size=14, row_h=Inches(0.46), hilite=[1])
line(s, L, Inches(4.4), Inches(7.2))
text(s, L, Inches(4.75), Inches(7.2), Inches(2.0), [
    [("MRR orders identifiers inside a fixed buffer. A bitmap scan "
      "materialises the whole qualifying set before touching the heap, so it "
      "has no such bound.", {"size": 15})],
    [("", {"size": 7})],
    [("At its best, on the one query class where it helps, MRR recovers about "
      "a third of one query's runtime against a cross-system gap of 15.4x.",
      {"size": 15, "bold": True})],
], spacing=1.35)
text(s, Inches(8.7), Inches(2.3), Inches(3.7), Inches(3.4), [
    [("A fourth backfiring practice", {"size": 14, "bold": True, "color": RED})],
    [("", {"size": 7})],
    [("\"Enable MRR\" is reasonable advice on the same grounds as the other "
      "three: the mechanism is real and the reasoning sound. At the shipped "
      "buffer size it makes the workload slower, and the second setting that "
      "would help is not mentioned alongside it.",
      {"size": 13.5, "color": MUTED})],
], spacing=1.32)
pagenum(s, 16)

# ================================================================ 17 BRIN
s = slide(dark=True, notes=(
    "Turning to PostgreSQL's own failure, which is the largest single effect "
    "in the study. Adding a BRIN index alongside an existing B-tree on the "
    "same column, because BRIN is small and people treat it as free. The "
    "bad-choice rate goes from 2.3 percent to 73.3. On a controlled test "
    "where the only change was whether the BRIN index existed, one query went "
    "from a third of a millisecond to sixty six."))
eyebrow(s, "Result 7  /  PostgreSQL-specific")
heading(s, "Adding an index makes the plan worse", color=ONDARK)
bignum(s, "194x", "slower on a controlled test where the only change\n"
                  "was that a BRIN index existed alongside the B-tree",
       L, Inches(2.4), Inches(6.1), size=96, color=RED_D, lab_color=DIMDK)
text(s, L, Inches(5.05), Inches(6.1), Inches(0.6),
     "0.34 ms  becomes  65.95 ms", size=22, font=SERIF, color=ONDARK,
     spacing=1.2)
grid(s, Inches(7.3), Inches(2.5),
     ["Table size", "With BRIN", "Without", "Worst case"],
     [["1,000,000", "73.3%", "2.3%", "18.8x"],
      ["10,000,000", "30.1%", "32.3%", "2.6x"]],
     [Inches(1.9), Inches(1.25), Inches(1.05), Inches(1.15)],
     aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 3,
     cell_size=14, row_h=Inches(0.5))
text(s, Inches(7.4), Inches(4.35), Inches(4.9), Inches(2.2), [
    [("MariaDB cannot form this plan at all, which is how we know it is a "
      "costing defect rather than something inherent to having two usable "
      "indexes.", {"size": 15, "color": ONDARK})],
    [("", {"size": 8})],
    [("Confined to tables that still fit in the buffer pool, which is exactly "
      "the size at which BRIN is recommended as a free addition.",
      {"size": 14, "color": DIMDK})],
], spacing=1.35)
pagenum(s, 17)

# ================================================================ 18 CAUSE
s = slide(notes=(
    "I traced that to a specific line rather than leaving it as an "
    "observation. When two indexes can answer the same predicate, PostgreSQL "
    "keeps one and discards the other, choosing whichever is cheaper to scan. "
    "But that cost is the cost of reading the index alone; it does not "
    "include the heap fetches the query then performs. A BRIN index is small "
    "by design, so it wins that comparison almost always, and the index being "
    "discarded is the one that would have run two hundred times faster."))
eyebrow(s, "Result 7, the mechanism")
heading(s, "Located in the source, not inferred")
text(s, L, Inches(2.05), Inches(6.3), Inches(0.5),
     "choose_bitmap_and()   in   optimizer/path/indxpath.c",
     size=15, font=MONO, color=BLUE, spacing=1.2)
text(s, L, Inches(2.7), Inches(6.3), Inches(2.6), [
    [("When two indexes can answer the same predicate, PostgreSQL keeps one "
      "and discards the other, choosing whichever is cheaper to scan.",
      {"size": 15.5})],
    [("", {"size": 7})],
    [("That cost is indextotalcost, the cost of reading the index alone. It "
      "excludes the heap fetches the query must then perform.",
      {"size": 15.5, "color": MUTED})],
    [("", {"size": 7})],
    [("A BRIN index is small by design, so it wins that comparison almost "
      "always, and the index thrown away is the faster one.",
      {"size": 15.5, "bold": True})],
], spacing=1.38)
grid(s, Inches(7.4), Inches(2.6),
     ["Index", "Cost used to decide", "Actual query time"],
     [["BRIN", "12.13   (kept)", "65.95 ms"],
      ["B-tree", "213.35   (discarded)", "0.34 ms"]],
     [Inches(1.3), Inches(2.3), Inches(1.75)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT],
     cell_size=14, row_h=Inches(0.52))
text(s, Inches(7.4), Inches(4.5), Inches(5.0), Inches(1.0),
     "The index that costs 17x less to read is the one that runs 194x slower.",
     size=15, color=RED, bold=True, spacing=1.3)
caption(s, "A design consequence of comparing index scan costs in isolation, which is why no setting removes it.")
pagenum(s, 18)

# ================================================================ 19 BOUNDARY
s = slide(notes=(
    "And the effect has a boundary that is worth reporting carefully, because "
    "it is where a two-point study would have gone wrong. The problem occurs "
    "while the table fits in the buffer pool, as expected. But read the two "
    "right-hand columns separately. How often it goes wrong collapses at the "
    "memory boundary, from 71 percent to 38. How badly it goes wrong doubles "
    "at that same boundary and stays high to three million rows. Testing only "
    "the smallest and largest sizes would have concluded the problem simply "
    "disappears as tables grow. It does not; the worst damage lands on "
    "databases sized two to three times their memory."))
eyebrow(s, "Result 8  /  reporting boundaries")
heading(s, "Frequency and severity move in opposite directions")
grid(s, L, Inches(1.92),
     ["Table size", "Heap", "In pool", "Chose badly", "Worst case"],
     [["1,000,000", "112 MB", "yes", "71.4%", "18.8x"],
      ["1,250,000", "140 MB", "no", "37.9%", "38.8x"],
      ["1,500,000", "168 MB", "no", "31.0%", "40.3x"],
      ["2,000,000", "223 MB", "no", "25.8%", "38.3x"],
      ["3,000,000", "335 MB", "no", "30.3%", "41.5x"],
      ["5,000,000", "558 MB", "no", "36.4%", "7.0x"],
      ["10,000,000", "1,116 MB", "no", "38.2%", "1.9x"]],
     [Inches(1.75), Inches(1.15), Inches(1.05), Inches(1.2), Inches(1.15)],
     aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT,
             PP_ALIGN.RIGHT],
     row_h=Inches(0.4), cell_size=13, header_h=Inches(0.56), hilite=[0, 4])
figure(s, "fig7_scale_transition.png", Inches(7.35), Inches(2.05), Inches(5.05), Inches(2.5))
text(s, Inches(7.35), Inches(4.75), Inches(5.05), Inches(1.9), [
    [("A two-point scale test would have reversed this conclusion.",
      {"size": 14.5, "bold": True, "color": RED})],
    [("", {"size": 6})],
    [("Frequency collapses at the memory boundary; severity doubles there and "
      "only subsides much later. The worst damage lands on databases sized "
      "two to three times their memory.", {"size": 13.5, "color": MUTED})],
], spacing=1.3)
pagenum(s, 19)

# ================================================================ 20 SCORECARD
s = slide(notes=(
    "Pulling the four practices together. Each one is sound in the case it "
    "was written for. Each one, under conditions I can now name, makes things "
    "worse. And in three of the four the advice improves the thing it "
    "promises to improve. Extended statistics really do fix the estimate, by "
    "a factor of 108. The plan still gets worse. That is the pattern the "
    "paper is really about: the advice targets a proxy, and the proxy is not "
    "what determines the outcome."))
eyebrow(s, "The pattern")
heading(s, "Four practices, each sound, each with an unstated boundary")
line(s, L, Inches(2.05), CW)
rows = [
    ("Extended statistics", "for correlated columns",
     "Fixes the estimate 108x", "Query 2.7x slower"),
    ("Lower random_page_cost", "because the disk is an SSD",
     "Default already within 6% of optimal", "Lowering it 2.3x worse"),
    ("Add a BRIN index", "it is small and cheap",
     "Bad choices 2.3% to 73.3%", "Worst case 194x"),
    ("Enable Multi-Range Read", "MariaDB, ships switched off",
     "Mechanism genuinely engages", "Workload 7% slower"),
]
y = Inches(2.45)
for name, why, does, costs in rows:
    text(s, L, y, Inches(3.5), Inches(0.8), [
        [(name, {"size": 15.5, "bold": True})],
        [(why, {"size": 12.5, "color": MUTED})],
    ], spacing=1.3)
    text(s, Inches(4.7), y + Inches(0.06), Inches(3.6), Inches(0.6), does,
         size=14, color=BLUE, spacing=1.25)
    text(s, Inches(8.6), y + Inches(0.06), Inches(3.8), Inches(0.6), costs,
         size=14, color=RED, bold=True, spacing=1.25)
    y += Inches(1.02)
    if y < Inches(6.3):
        line(s, L, y - Inches(0.16), CW)
text(s, L, Inches(6.55), CW, Inches(0.5),
     "In three of four the advice delivers exactly what it promises, and the "
     "plan gets worse anyway.", size=14.5, bold=True, spacing=1.25)
pagenum(s, 20)

# ================================================================ 21 VALIDITY
s = slide(notes=(
    "On trustworthiness, briefly, because this is where a reviewer presses. "
    "Nineteen automated checks run before any measurement is accepted, "
    "including validating the generated data against PostgreSQL's own "
    "internal statistics, so the generator is checked against something "
    "outside my control. Dispersion is reported rather than hidden. And three "
    "of my own errors were found and are reported rather than removed. The "
    "one I would highlight is the third: I had claimed the correlated-column "
    "failure was caused by the plan shape that merges two indexes. I tested "
    "that by switching the plan off, and the error stayed the same, 55.59 "
    "against 55.84. So the merge is where the error becomes visible, not "
    "where it comes from. I only found the real cause by trying to refute "
    "myself."))
eyebrow(s, "Validity")
heading(s, "Including one explanation I had to withdraw")
line(s, L, Inches(2.05), CW)
text(s, L, Inches(2.45), Inches(3.6), Inches(3.6), [
    [("19", {"size": 44, "font": SERIF, "color": BLUE})],
    [("automated checks before any measurement is accepted",
      {"size": 13, "bold": True})],
    [("", {"size": 6})],
    [("Generated data validated against PostgreSQL's own pg_stats, so the "
      "generator is checked against something outside my control.",
      {"size": 13, "color": MUTED})],
], spacing=1.3)
text(s, Inches(4.9), Inches(2.45), Inches(3.4), Inches(3.6), [
    [("1.79%", {"size": 44, "font": SERIF, "color": BLUE})],
    [("median coefficient of variation above the 1 ms floor",
      {"size": 13, "bold": True})],
    [("", {"size": 6})],
    [("24.47% below it, which is why that floor exists. Outlying runs are "
      "retained, not deleted; the median absorbs up to two in six.",
      {"size": 13, "color": MUTED})],
], spacing=1.3)
text(s, Inches(9.0), Inches(2.45), Inches(3.4), Inches(3.6), [
    [("3", {"size": 44, "font": SERIF, "color": RED})],
    [("of my own errors found, corrected and reported",
      {"size": 13, "bold": True})],
    [("", {"size": 6})],
    [("Including a claim I refuted myself: disabling the merge plan left the "
      "error at 55.84 against 55.59, so the merge shows the error rather than "
      "causing it.", {"size": 13, "color": MUTED})],
], spacing=1.3)
caption(s, "All 23,480 raw measurements, seeded generators and the analysis notebook are public and re-run with one command.")
pagenum(s, 21)

# ================================================================ 22 THESIS
s = slide(notes=(
    "What the study actually identifies, in four claims. First, for this "
    "decision the cause is the cost model, not the statistics, and 98 percent "
    "of bad choices with correct estimates is the evidence. Second, tuning "
    "advice needs its conditions attached, because all four practices are "
    "sound in one regime and harmful in another. Third, an optimiser cannot "
    "be summarised by one number, because frequency and severity move in "
    "opposite directions across scale and across these two engines. Fourth, "
    "and this is what the second database bought, per-column composition "
    "under correlation is a shared design weakness rather than a bug in "
    "either product."))
eyebrow(s, "What the study identifies")
heading(s, "Four claims")
line(s, L, Inches(2.05), CW)
concl = [
    ("The cause is the cost model, not the statistics",
     "98.2% of bad choices carried correct row estimates. Better statistics "
     "cannot repair a comparison that omits heap work, so the fix belongs in costing."),
    ("Tuning advice is incomplete without its conditions",
     "All four practices tested are sound in the case they were written for "
     "and harmful outside it. None ships with that boundary stated."),
    ("One number cannot describe an optimiser",
     "How often it errs and how badly it errs move in opposite directions "
     "across scale, and in opposite directions between these two engines."),
    ("Per-column composition is a shared weakness, not a bug",
     "Two independently written engines are accurate only while a single "
     "index covers both columns, and both fail once none does."),
]
y = Inches(2.4)
for i, (h, b) in enumerate(concl, 1):
    text(s, L, y, Inches(0.5), Inches(0.4), str(i), size=15, font=SERIF,
         bold=True, color=RED, spacing=1.0)
    text(s, L + Inches(0.55), y - Inches(0.02), Inches(11.0), Inches(1.0), [
        [(h, {"size": 16.5, "bold": True})],
        [(b, {"size": 14, "color": MUTED})],
    ], spacing=1.3)
    y += Inches(1.12)
pagenum(s, 22)

# ================================================================ 23 CLOSE
s = slide(dark=True, notes=(
    "That is the work. Two engines, seven scales, twenty three thousand "
    "measurements, everything public and reproducible with one command. Happy "
    "to take questions, and if you want any number in the deck verified I can "
    "show it coming out of the raw data."))
line(s, L, Inches(2.0), Inches(1.4), color=RED, weight=2.5)
text(s, L, Inches(2.35), Inches(9.5), Inches(1.4),
     "Thank you. Questions?", size=44, font=SERIF, color=ONDARK, spacing=1.1)
line(s, L, Inches(3.75), CW, color=RGBColor(0x33, 0x38, 0x3D))
text(s, L, Inches(4.1), Inches(5.6), Inches(2.2), [
    [("Everything is reproducible", {"size": 14, "bold": True, "color": ONDARK})],
    [("", {"size": 6})],
    [("23,480 raw measurements with every repetition retained, seeded "
      "generators, the analysis notebook and both write-ups.",
      {"size": 14, "color": DIMDK})],
], spacing=1.35)
text(s, Inches(7.2), Inches(4.1), Inches(5.2), Inches(2.2), [
    [("Md. Imtiaj Alam Sajin", {"size": 15, "bold": True, "color": ONDARK})],
    [("26-94090-2", {"size": 13, "color": DIMDK})],
    [("Supervisor: Dr. Ashraf Uddin", {"size": 13, "color": DIMDK})],
], spacing=1.4)
text(s, L, H - Inches(0.85), CW, Inches(0.4),
     "github.com/Imtiaj-Sajin/Reseach-on-Database-Architecture",
     size=11.5, font=MONO, color=DIMDK, spacing=1.2)

prs.save(OUT)
print("wrote", OUT)
print("slides:", len(prs.slides._sldIdLst))
